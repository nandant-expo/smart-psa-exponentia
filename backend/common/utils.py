import os
import re
import sys
import fitz
import nltk
import pytz
import uuid
import time
import json
import PyPDF2
import random
import openai
import logging
import tiktoken
import requests
import threading
import numpy as np
from PIL import Image
import pandas as pd
from PIL import Image
from io import BytesIO
from retry import retry
import concurrent.futures
import aspose.words as aw
import aspose.cells  as ac
from threading import Lock
from fuzzywuzzy import fuzz
from itertools import repeat
from whoosh.qparser import *
from datetime import datetime,timedelta
from pptx import Presentation
import azure.cognitiveservices.speech as speechsdk
from azure.cognitiveservices.speech import AudioConfig, SpeechConfig, SpeechRecognizer, AutoDetectSourceLanguageConfig
import aspose.slides as slides
import aspose.pydrawing as drawing
from bson import json_util,ObjectId
from detect_delimiter import detect
from common.database import Database
from cachetools import TTLCache, cached
from pdf2image import convert_from_bytes
from transformers import GPT2TokenizerFast
from collections import Counter,defaultdict
from pptx.enum.shapes import MSO_SHAPE_TYPE
from azure.search.documents import SearchClient
from azure.storage.blob import BlobServiceClient
from requests.exceptions import RequestException
from azure.search.documents.indexes.models import *
from azure.search.documents.models import QueryType
from azure.core.credentials import AzureKeyCredential
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.search.documents.indexes import SearchIndexClient 
from passlib.context import CryptContext
from common.keyvault_connection import get_conn
from difflib import SequenceMatcher

ist_timezone = pytz.timezone('Asia/Kolkata')


nltk.download('stopwords')
from nltk.corpus import stopwords

client=get_conn()
storage_connection_string = client.get_secret('STORAGE-CONNECTION-STRING').value
storage_account = client.get_secret("STORAGE-ACCOUNT").value

endpoint = client.get_secret("AZURE-DOC-INTELLIGENCE-ENDPOINT").value
credential = AzureKeyCredential("AZURE-DOC-INTELLIGENCE-KEY")
document_analysis_client = DocumentAnalysisClient(endpoint, credential)

tokenizer = GPT2TokenizerFast.from_pretrained("gpt2")
openai.api_type = "azure"
openai.api_version = "2023-05-15"
openai_api_base = [client.get_secret("AZURE-OPENAI-ENDPOINT-US").value,client.get_secret("AZURE-OPENAI-ENDPOINT-EUROPE").value]
openai_key = [client.get_secret('OPENAI-API-KEY-US').value,client.get_secret('OPENAI-API-KEY-EUROPE').value]

AZURE_SEARCH_SERVICE = client.get_secret("AZURE-SEARCH-SERVICE").value
AZURE_SEARCH_KEY = client.get_secret("AZURE-SEARCH-KEY").value
searchindexclient = SearchIndexClient(endpoint=f"https://{AZURE_SEARCH_SERVICE}.search.windows.net",
    credential=AzureKeyCredential(AZURE_SEARCH_KEY))

password_context = CryptContext(schemes=["bcrypt"], deprecated="auto")

workspace_url = client.get_secret("DATABRICKS-URL").value
token = client.get_secret("DATABRICKS-PAT").value

base_url = "https://graph.microsoft.com/v1.0/sites"
kb_index = "{}-gptxponent-chat"
repo_index = "{}-gptxponent-repo"
search_index = "{}-file-index"
chunk_index = "{}-chunk-index"

search_master_table="search_document_collection"
search_page_table="search_page_collection"
PAT_TOKEN = token
DATABRICKS_URL = workspace_url
CATALOG = "databricks_gptexponent_catalogue"

def get_documents_list_id(access_token,site_id):
    url = f"{base_url}/{site_id}/lists"
    ids = []
    payload = ""
    headers = {
    'Content-Type': 'application/json',
    'Authorization': f'Bearer {access_token}'
    }

    response = requests.request("GET", url, headers=headers, data=payload)
    response = response.json()
    if 'value' in response:
        ids.extend([item['id'] for item in response['value'] if item['displayName']=='Documents'])
    return ids

def sync_file_access(user_id,access_token,workspace):

    db =Database('pipelineSettings',workspace)

    pipeline = [{"$group": {"_id": "$siteid","foldername": {"$addToSet": "$foldername"}}}]
    
    status,records = db.fetch_aggregate(pipeline)
    
    item_list = []
    if status:
        # print(records)
        for record in records:
            site_id = record['_id']
            list_ids = get_documents_list_id(access_token,site_id)
            for list_id in list_ids:
                url = f"{base_url}/{site_id}/lists/{list_id}/items"
                payload = {}
                headers = {
                    'Authorization': f'Bearer {access_token}'
                }
                # response = requests.request("GET", url, headers=headers, data=payload)
                # response = response.json()
                
                data=",".join(record['foldername']).replace(' ', '%20')
                # print(data)
                # print(sizes)
                while True:
                    response = requests.request("GET", url, headers=headers, data=payload)
                    response = response.json()
                    if 'value' in response:
                        for items in response['value']:
                            match = SequenceMatcher(None, data, items['webUrl'],autojunk=False).find_longest_match(0,len(data),0,len(items['webUrl']))
                            # items['webUrl']
                            try:
                                index=record['foldername'].index(data[match.a:match.a + match.size].replace('%20', ' '))
                                # print(index)
                            except ValueError:
                                # print(match)
                                # print(items['webUrl'])
                                # print(data[match.a:match.a + match.size])
                                index=-1
                            
                            if index!=-1 and items['webUrl'].split(".")[-1].lower() in ['pptx','ppt','pdf','docx','doc','xls','png','jpeg','jpg','xlsx','mp4','mp3','wav','amr','m4a']:
                                # print(data[match.a:match.a + match.size])
                                item_list.append(items['webUrl'].split("Shared%20Documents")[-1])
                    if '@odata.nextLink' not in response.keys():
                        break
                    else:
                        url = response['@odata.nextLink']
            # break
    
    db = Database('DocumentCollection',workspace)
    # print(item_list)
    status,records = db.fetch_all_records({"item_name": {"$nin": item_list},"file_source":"sharepoint"},{'_id':1,"display_name":1})
    not_allowed_ids = []
    not_allowed_names = []
    if status:
        not_allowed_ids=list(map(lambda item: item['_id'], records))
        not_allowed_names=list(map(lambda item: item['display_name'], records))
        
    db = Database('UserActivity',workspace)
    status,message = db.update_record({"user_id":user_id},{'listing_filter':not_allowed_ids,'chat_filter':not_allowed_names},{'$pull':{'file_history':{'file_id':{'$in':not_allowed_ids}}}})
    if status:
        return True,message
    return False,message

def check_credentials(username, password, loginType,database_name):
    try:
        db = Database('users',database_name)
        query = [{'$match': {'email': username}},
        {'$project':{'_id': 0, 'user_id': '$_id', 'role': 1, 'password' :{'$ifNull': [ "$password", None ] },'email':1,'displayName':1,'workspace':1}}]
        
        status,result = db.fetch_aggregate(query)
        if status:
            result = result[0]
            if len(result) > 0:
                db = Database('UserManagement',result['workspace'])
                query = [{'$match': {'username': username}},{'$project':{'_id': 0, 'user_id': '$_id'}}]
                if loginType==None: 
                    if password_context.verify(password, result['password']):
                        status_data,result_data = db.fetch_aggregate(query)
                        if status_data and len(result_data)>0:
                            return True, (result['role'],json_util.dumps(result_data[0]['user_id']),result['email'],result['displayName'],True,result['workspace'])
                if password == result['password'] and loginType=='AD':
                    status_data,result_data = db.fetch_aggregate(query)
                    if status_data:
                        return True, (result['role'],json_util.dumps(result_data[0]['user_id']),result['email'],result['displayName'],True,result['workspace'])
        else:
            if password==None and loginType=='AD':
                return True, (None,json_util.dumps(None),None,None,False,None)
        print("Incorrect username or password")
        return False, 'Incorrect username or password'    
    except Exception as error:
        print(f"utility function - check_credentials - Error : {error}")
        return False, "Failed to validate credentials"

def blob_connection():
    try:
        client = BlobServiceClient.from_connection_string(storage_connection_string)
        # print("Blob Storage Client has been created successfully!")
        return True, client
    except Exception as error:
        print(f"utility function - blob_connection - Error : {error}")
        return False, error

def delete_blob(path):
    try:
        client_staus, blob_client = blob_connection()
        if client_staus:
            container = path.split("/")[3]
            path = path.split(container + "/")[-1]
            blob_data = blob_client.get_blob_client(
                container, path)
            blob_data.delete_blob()
            print("File has been succeessfully Deleted.")
            return True, None
        else:
            print(
                f"Failed to delete the file - Blob connection failed. {path.split('/')[-1]} - error {blob_client}")
        return False, 'Failed to delete the file.'
    except Exception as error:
        print(
            f"ERROR: delete_blob -> filename: {path.split('/')[-1]} -> Error: {error}")

@cached(cache = TTLCache(maxsize=10, ttl=30), lock=Lock())
def download_blob(path):
    try:
        client_staus, blob_client = blob_connection()
        filename = str(path.split("/")[-1])
        container = path.split("/")[3]
        print(container)
        print(path.split(f"{container}/")[-1])
        if client_staus:
            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Blob client successfully created")
            blob = blob_client.get_blob_client(
                container, path.split(f"{container}/")[-1])
            
            blob_data = blob.download_blob()
            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Blob file downloaded")

            created_on = blob.get_blob_properties().creation_time.strftime('%Y-%m-%d %H:%M:%S')
            blob_data = blob_data.readall()
            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Blob file metadata fetched")

            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: File has been succeessfully loaded.")
        else:
            print(
                f"Failed to download the file - Blob connection failed. {filename} - error {blob_client}")
        return True, {"file" : blob_data, "created_on" : created_on}
    except Exception as error:
        print(
            f"ERROR: download_blob -> filename: {filename} -> Error: {error}")
        return False, error
    
def upload_image_thumbnail(file_name,type,image,container=False):
    main_dir = f"slide-data/file_{type}" if not container else f"{container}/file_{type}"
    storage_uri = storage_account
    image_blob_path = os.path.join(main_dir, file_name)
    container_name = client.get_secret("CONTAINER-NAME").value
    
    blob_service_client = BlobServiceClient.from_connection_string(storage_connection_string)
    container_client=blob_service_client.get_container_client(container_name)
    blob_client = container_client.get_blob_client(image_blob_path)

    blob_client.upload_blob(image.read(),overwrite=True)
    fileurl = os.path.join(storage_uri, container_name, image_blob_path)
    return fileurl

def get_user_settings(col=False,workspace='smartsearch-exponentia-dev'):
    try:
        cols = {}
        if col:
            cols.update(col)
        status,result = Database('UserConfiguration',workspace).fetch_all_records(cols=cols)
        if status:
            return status,result[0]
        return 
    except Exception as error:
        print(f'Failed:get_user_settings.Error: {error}')

def get_one_record(id,workspace):
    try:
        db = Database('DocumentCollection',workspace)
        query = {"_id":id}
        cols = {
            '_id' : 1,
            'file_name' : 1,
            'display_name':1,
            'file_type':1,
            'file_path':1,
            'file_title' : 1,
            'file_summary':1,
            'tags' : 1,
            'file_entities' : 1,
            'file_dictionary' : 1,
            'file_preview' : 1 
            }
        status,result =  db.fetch_one_record(query,cols)
        if status:
            if len(result) > 0:
                status, suggestion = get_user_settings(col={'_id':0,'suggestions':1},workspace=workspace)
                result.update(suggestion)
                return True, result
            return False, "No records found"
        return False, "Failed to fetch records"
    except Exception as error:
        print(f"utility function - get_one_record - Error : {error}")
        return False, "Failed to fetch records."
    
def edit_ppt_data(id,data,workspace):
    try:
        db_central=Database('workspace',"central")
        status_central,result_central = db_central.fetch_one_record({"workspace_name" : workspace},{"upload_job_id":1})
        if status_central and "upload_job_id" in result_central:
            utc_now = datetime.utcnow()
            utc_timestamp_str = utc_now.strftime('%Y-%m-%d %H:%M:%S')
            db = Database('DocumentCollection',workspace)
            status,result = db.fetch_one_record({"_id":id},{"tags":1,"word_cloud":1,'user_tags':1})
            add_data = {}
            deleted_items = [item for item in result['tags'] if item not in data['tags']]
            add_data['word_cloud'] = result['word_cloud']
            extra_items = [item for item in data['tags'] if item not in add_data['word_cloud']]
            add_data["word_cloud"][:0] = extra_items
            add_data['word_cloud'] = sorted(
                [max(
                    [item for item in add_data['word_cloud'] if item['word'] == word],
                    key=lambda x: x['weight']
                ) for word in set(item['word'] for item in add_data['word_cloud'])],
                key=lambda x: x['weight'],
                reverse=True
            )
            add_data['user_tags'] = []
            if 'user_tags' in result.keys():
                add_data['user_tags'] = result['user_tags']
            add_data['user_tags'].extend(extra_items)
            # add_data['edit_flag'] = 1
            add_data['form_modified_on'] = utc_timestamp_str
            if deleted_items:
                for item in deleted_items:
                    if item in add_data['word_cloud']:
                        add_data['word_cloud'].remove(item)
                    if item in add_data['user_tags']:
                        add_data['user_tags'].remove(item)
            data.update(add_data)
            status,message = db.update_one_record(id,data)
            if status:
                print("PPT data has been updated.")
                wokspaceurl=DATABRICKS_URL
                databricks_api_url = f"{wokspaceurl}/api/2.1/jobs/run-now"
                token = PAT_TOKEN
                
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {token}"
                }
                payload = {
                    "job_id": result_central['upload_job_id'],
                    "job_parameters": {"type":"update","doc_id":str(id), "workspace":workspace}
                }

                logging.info("Sending request to Databricks...")
                response = requests.post(databricks_api_url, headers=headers, json=payload)

                if response.status_code==200:
                    logging.info("Databricks Job Triggered")
                else:
                    logging.info(f"Something went wrong during job trigger: {response.text}")
                return True, message
            return False, message
        return False,"No Upload Job Available"
    except Exception as error:
        print(f'Failed ppt_methods->update_ppt_data. Error:{error}.')
        return False, "Failed to edit records."
    
def delete_data(id,workspace):
    try:
        # print(id)
        db_central=Database('workspace',"central")
        status_central,result_central = db_central.fetch_one_record({"workspace_name" : workspace},{"upload_job_id":1})
        if status_central and "upload_job_id" in result_central:
            db = Database('DocumentCollection',workspace)
            print("Retrieving the record to be deleted.")
            status,result = db.fetch_one_record({"_id" : id})
            if status:
                print("Retrieved the record to be deleted.")
                
                wokspaceurl=DATABRICKS_URL
                databricks_api_url = f"{wokspaceurl}/api/2.1/jobs/run-now"
                token = PAT_TOKEN
                
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {token}"
                }
                payload = {
                    "job_id": result_central['upload_job_id'],
                    "job_parameters": {"type":"delete","doc_id":str(id), "workspace":workspace}
                }

                logging.info("Sending request to Databricks...")
                response = requests.post(databricks_api_url, headers=headers, json=payload)

                if response.status_code==200:
                    logging.info("Databricks Job Triggered")
                else:
                    logging.info(f"Something went wrong during job trigger: {response.text}")
                if result['file_source'] == 'upload':
                    print(result['file_path'])
                    status,message = delete_blob(result['file_path']) if 'file_path' in result.keys() else True,None
                i=0
                if 'file_preview' in result.keys():
                    while status and i<len(result['file_preview']):
                        status,message = delete_blob(result['file_preview'][i]['thumbnail_path'])
                        print(message)
                        status,message = delete_blob(result['file_preview'][i]['image_path'])
                        print(message)
                        i+=1
                result = message
                print("Deleted the file from blob storage.")
                status,result = db.delete_single_record(id)
                if status:
                    print("Deleted the record from database.")
                    return True, None
            return False, result
        return False,"No Upload Job Exists"
    except Exception as error:
        print(f"utility function - delete_data - Error : {error}")
        return False,"Failed to delete file and record."

def fetch_records(user_id,workspace):
    try:
        exclude_ids = []
        db = Database('UserActivity',workspace)
        status,record = db.fetch_one_record({'user_id':user_id},{'listing_filter':1})
        if status:
            if 'listing_filter' in record:
                exclude_ids = record['listing_filter']
        db = Database('DocumentCollection',workspace)
        print("Retrieving the records.")
        filter = {"edit_flag":1,"_id": {"$nin": exclude_ids}}
        query = {
            '_id':1,
            'file_name':1,
            'display_name':1,
            'file_type':1,
            'file_title':1,
            'file_source':1,
            'tags':1,
            'uploaded_on':1,
            'file_path':1,
            'file_preview' : 1,
            'status':1
            }
        status,result = db.fetch_all_records(filter=filter,cols = query)
        if status:
            print("Sucessfully retrieved records.")
            return True,result
        print("Failed to retrieve records")
        if result == "No records found.":
            return True,[]
        else:
            return False, result
    except Exception as error:
        print(f"utility function - fetch_records - Error : {error}")
        return False,"Failed to fetch records."

def fetch_sharepoint_records(workspace):
    try:
        db = Database('PidiliteDocuments',workspace)
        print("Retrieving the records.")
        filter = {}
        query = {
            '_id':1,
            'file_name':1,
            'display_name':1,
            'file_type':1,
            'file_title':1,
            'tags':1,
            'created_on':1,
            'modified_on':1,
            'file_preview' : 1 
            }
        status,result = db.fetch_all_records(filter=filter,cols = query)
        if status:
            print("Sucessfully retrieved records.")
            return True,result
        print("Failed to retrieve records")
        return False, result
    except Exception as error:
        print(f"utility function - fetch_records - Error : {error}")
        return False,"Failed to fetch records."
    
def remove_from_index(id,workspace):
    db_central=Database("workspace","central")
    db=Database("DocumentCollection",workspace)
    job_id = None
    job_id = ''
    if job_id:
        workspaceurl=DATABRICKS_URL
        databricks_api_url = f"{workspaceurl}/api/2.1/jobs/run-now"
        token = PAT_TOKEN
        print(f'ID: {id}, {str(id)}, {json_util.dumps(id)}')
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {token}"
        }
        payload = {
            "job_id": job_id,
            "job_parameters": {"type":"delete","doc_id":str(id), "workspace":workspace}
        }
 
        logging.info("Sending request to Databricks...")
        response = requests.post(databricks_api_url, headers=headers, json=payload)
        print('Response: ', response.json())
        if response.status_code==200:
            print('Deleted file from index successfully')
            return True
        return False

@retry(Exception, tries=10, delay=1)
def keywords_extraction(input_text):
    try:
        def chunk_text_by_tokens(text, chunk_size):
            tokens = tokenizer.tokenize(text)
            chunks = []

            for i in range(0, len(tokens), chunk_size):
                chunk = tokens[i:i + chunk_size]
                # print(len(chunk))
                chunks.append(chunk)

            return chunks

        def extract_keywords_weightages(output):
            # Unwanted characters removal (,;)
            output = re.sub(r"(?<!\d)[.;]|[.;](?!\d)", " ", output)

            # Extract keyword:weightage
            pairs = re.findall(r"([\w\s]+):?\s*([\d.]*)", output)

            keywords_weightages = []
            for pair in pairs:
                keyword, weightage = pair
                weightage = float(weightage) if weightage else 0
                keywords_weightages.append((keyword.strip(), weightage))

            return keywords_weightages
        
        def merge_responses(responses):
            merged_dict = {}

            for response in responses:
                keywords_weightages = extract_keywords_weightages(response)
                
                for keyword, weightage in keywords_weightages:
                    
                    if keyword in merged_dict and merged_dict[keyword] < weightage:
                        merged_dict[keyword] = weightage
                    
                    elif keyword not in merged_dict:
                        merged_dict[keyword] = weightage

            
            merged_list = [(keyword, weightage) for keyword, weightage in merged_dict.items()]
            return merged_list
        
        @retry(Exception, tries=10, delay=1)
        def openai_gpt_model(prompt):
            try:
                rindex = random.randint(0,1)
                openai.azure_endpoint = openai_api_base[rindex]
                openai.api_key = openai_key[rindex]
                messages = [{"role": "user", "content": f"Analyze the following text and consistently identify the most important keywords on multiple runs. Ensure that the keywords are no longer than three words each. Then, assign a weight to each keyword on a scale of 0 to 1, representing its importance within the text. Use the format 'keyword:weight' and ensure the format is consistent across multiple runs and does not have unwanted characters in the output response: {prompt} Consistent keywords and weights in the desired format:"}]
                response = openai.ChatCompletion.create(
                    engine="chatgpt",
                    messages=messages,
                    temperature=0,
                    request_timeout=20
                )
                text = [response.choices[0].message["content"].strip() for i in range(2)]
                output = merge_responses(text)
                return output
            except Exception as error:
                print(f"Failed->keywords_extraction->openai_gpt_model. Error {error}")
                return ""


        
        chunk_size = 3000
        chunks = chunk_text_by_tokens(input_text, chunk_size)
        chunk_texts = []
        for i, chunk in enumerate(chunks):
            readable_chunk = tokenizer.convert_tokens_to_string(chunk)
            chunk_texts.append(readable_chunk)
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            # Submit the processing function for each chunk and collect the futures
            futures = [executor.submit(openai_gpt_model, chunk) for chunk in chunk_texts]

            # Iterate over the completed futures and flatten the results
            keyword_weights = [tup for sublist in [future.result() for future in concurrent.futures.as_completed(futures)] for tup in sublist if tup]
        
        combined_keywords_weightage = defaultdict(float)
        keyword_counts = defaultdict(int)

        for keyword, weightage in keyword_weights:
            combined_keywords_weightage[keyword] += weightage
            keyword_counts[keyword] += 1

        average_keywords_weightage = {
            keyword: (combined_keywords_weightage[keyword] / keyword_counts[keyword])
            for keyword in combined_keywords_weightage
        }
        final_keyword_list = dict([(keyword,float("%.2f" % round(weightage, 2))) for keyword,weightage in average_keywords_weightage.items()])
        if len(final_keyword_list)>10:
            final_keyword_list = dict([(keyword,float("%.2f" % round(weightage, 2))) if weightage>0.2 and len(average_keywords_weightage.items())>10 and keyword else (keyword,float("%.2f" % round(weightage, 2))) for keyword,weightage in average_keywords_weightage.items() if keyword])

        return final_keyword_list
    except Exception as error:
            print(f"Failed->keywords_extraction. Error {error}")

def get_keywords(text):
    try:
        keywords = keywords_extraction(text)
        keywords = dict(sorted(keywords.items(), key=lambda x: x[1],reverse=True))
        keywords = list(keywords.items())
        word_cloud = [{"word":key.title(),"weight":value} for (key, value) in keywords[:50] if key and len(key)>2]
        tags = [{"word":key.title(),"weight":value} for (key, value) in keywords[:10] if key and len(key)>2]
        return tags,word_cloud
    except Exception as error:
        print(f'Failed nlp_methods->get_keywords. Error:{error}')

@retry(Exception, tries=5, delay=1)
def ner(prompt):
    try:
        rindex = random.randint(0,1)
        openai.api_base = openai_api_base[rindex]
        openai.api_key = openai_key[rindex]
        messages = [{"role": "user", "content": prompt}]
        response = openai.ChatCompletion.create(
            engine="chatgpt",
            messages=messages,
            temperature=0,
            request_timeout=20
        )
        text = response.choices[0].message["content"]
        if text:
            return text.strip()
        return ""   
    except Exception as error:
        print(f"Failed due to {error}")

# ner_couts = []
def extract_ner(text,workspace):
    try:
        status, ner_configuration = get_user_settings(col={'_id':0,'ner':1},workspace=workspace)
        if status:
            ner_setting = {category['category']:[{entity:round(random.uniform(0.5, 0.99),2)} for entity in category['entities']] for category in ner_configuration['ner']}
        ner_setting_no_output = dict([(category,[]) for category in ner_setting.keys()])
        prompt = f"""Identify the entities in the text according to the following categories: {", ".join([i for i in list(ner_setting.keys())])}.
                    For Example:
                    Output:
                    {ner_setting}
                    If no output found, provide output as following:
                    {ner_setting_no_output}
                    Text:{text}
                    Do not mention entities outside the given text data and I'm expecting entities as bigram and trigram.
                    Output:"""
        response = ner(prompt).splitlines()[0]
        if isinstance(response,dict):
            return response
        ner_output = eval(response)
        ner_output_final = {}
        if ner_output:
            # ner_couts.append(ner_output)
            for category, entities in ner_output.items():
                ner_output_final[category] = [name for entity in entities for name, score in entity.items() if score >= 0.75 and len(name.split())<=3]
        return ner_output_final
    except Exception as error:
        print(f"Error:{error}")
        return ner_setting_no_output


def ocr_img2text(image_stream):
    SUBSCRIPTION_KEY = client.get_secret("OCR-SUBSCRIPTION").value
    ENDPOINT = client.get_secret("OCR-ENDPOINT").value
    headers = {
        'Ocp-Apim-Subscription-Key': SUBSCRIPTION_KEY,
        'Content-Type': 'application/octet-stream'
    }
    params = {
        'language': 'en',
        'detectOrientation': 'true'
    }
    read_url = f"{ENDPOINT}vision/v3.2/read/analyze"

    try:
        response = requests.post(read_url, headers=headers, params=params, data=image_stream)
        
        operation_location = response.headers["Operation-Location"]
        
        analysis = {}
        while not "analyzeResult" in analysis:
            response_final = requests.get(operation_location, headers=headers)
            analysis = response_final.json()
            # time.sleep(1)

        text_lines = []
        for read_result in analysis['analyzeResult']['readResults']:
            # if not 'lines' in read_result:
            #     return True,''
            for line in read_result['lines']:
                text_lines.append(re.sub(r'([^a-zA-Z ]+?)','',line['text']))
        return True," ".join(text_lines)
    except Exception as error:
        logging(f'Failed: ocr_img2text. Error: {error}')
        return False, "Failed to perform OCR operation."
    
def get_PPTimage_text(shape):
    try:
        img = shape.image.blob
        text = ''
        # text = image_to_text(img)
        status,text = ocr_img2text(bytearray(img))
        if status:
            if text is not None:
                # return " ".join(text.splitlines())
                return text
        return ''
    except Exception as error:
        print(f'Failed ppt_methods->get_PPTimage_text. Error:{error}')
        return ''


@retry(Exception, tries=10, delay=1)
def summary_generation(input_text):
    try:
        def chunk_text_by_tokens(text, chunk_size):
            tokens = tokenizer.tokenize(text)
            chunks = []

            for i in range(0, len(tokens), chunk_size):
                chunk = tokens[i:i + chunk_size]
                # print(len(chunk))
                chunks.append(chunk)

            return chunks

        @retry(Exception, tries=10, delay=1)
        def summary_bullet_points(prompt):
            try:
                rindex = random.randint(0,1)
                openai.api_base = openai_api_base[rindex]
                openai.api_key = openai_key[rindex]
                messages = [{"role": "user", "content": f"Please provide a complete summary of the following text in 10 bullet points with consistent summary and concise information:\n\n{prompt}\n\nBullet Points:\n-"}]
                response = openai.ChatCompletion.create(
                    engine="chatgpt",
                    messages=messages,
                    temperature=0,
                    request_timeout=20
                )
                text = response.choices[0].message["content"]
                if text:
                    return text.strip()
                return ""   
            except Exception as error:
                print(f"Failed->summary_generation->summary_bullet_points. Error {error}")
                return ""

        @retry(Exception, tries=10, delay=1)
        def summary_gpt_model(prompt):
            try:
                rindex = random.randint(0,1)
                openai.api_base = openai_api_base[rindex]
                openai.api_key = openai_key[rindex]
                messages = [{"role": "user", "content": f"Please summarize the following text without mentioning any years:{prompt} Summary:"}]
                response = openai.ChatCompletion.create(
                    engine="chatgpt",
                    messages=messages,
                    temperature=0,
                    request_timeout=20
                )
                text = response.choices[0].message["content"]
                if text:
                    return text.strip()
                return ""
            except Exception as error:
                print(f"Failed->summary_generation->summary_gpt_model. Error {error}")
                # print(prompt)
                return ""

        chunk_size = 3000
        chunks = chunk_text_by_tokens(input_text, chunk_size)
        chunk_texts = []
        for i, chunk in enumerate(chunks):
            readable_chunk = tokenizer.convert_tokens_to_string(chunk)
            chunk_texts.append(readable_chunk)

        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            futures = [executor.submit(summary_gpt_model, chunk) for chunk in chunk_texts]

            # Iterate over the completed futures and flatten the results
            summaries = [tup for sublist in [future.result() for future in concurrent.futures.as_completed(futures)] for tup in sublist]
        summaries = "".join(summaries)
        return summary_bullet_points(summaries)
    except Exception as error:
        print(f"Failed->summary_generation. Error {error}")

def num_tokens_from_text(messages):
    encoding= tiktoken.get_encoding("cl100k_base")  
    num_tokens = 0
    num_tokens += len(encoding.encode(messages))
    return num_tokens

@retry(Exception, tries=10, delay=1)
def doc_summary_generation(input_text,doc_type):
    try:
        rindex = random.randint(0,1)
        openai.api_base = openai_api_base[rindex]
        openai.api_key = openai_key[rindex]
        summary_text = input_text
        summary_text = summary_text.split()
        summary_text_limit =''
        summary_tokens = num_tokens_from_text(summary_text_limit)
        response_tokens = 500
        token_limits = [3000,3000]
        while len(summary_text)>0:
            summary_text_limit = summary_text_limit + summary_text[0] + " "
            del summary_text[0] 
            summary_tokens = num_tokens_from_text(summary_text_limit)
            if summary_tokens + response_tokens >= token_limits[rindex]:
                break
        messages = [
            {"role":"system","content":"""You are a summary generator agent.\nprovide a complete summary of the following in atmost one paragraph with consistent summary and concise information.\n
            Example:\n
            Text:Harnessing sustainable sources for cleaner, greener power solutions.\n
            File Type:pdf\n\n
            Summary:PDF on clean, green power solutions from sustainable sources.\n\n\n
            Text:Harnessing sustainable sources for cleaner, greener power solutions.\n
            File Type: video\n\n
            Summary:Video on tapping sustainable sources for cleaner, greener power solutions.\n\n\n
            Text: Analyzing global environmental trends and their impact on business strategies.\n
            File Type: ppt\n\n
            Summary: Presentation (PPT) exploring the influence of global environmental trends on business strategies worldwide.\n\n\n
            Text: A comprehensive guide to implementing sustainable practices in corporate settings.\n
            File Type: doc\n\n
            Summary: Document providing a detailed roadmap for incorporating sustainable practices in corporate environments.\n\n\n
            Text: Cutting-edge innovations in renewable energy technologies for a sustainable future.\n
            File Type: img\n\n
            Summary: Image file showcasing the latest innovations in renewable energy technologies, contributing to a sustainable future.\n\n\n
            Text: Financial analysis of clean energy projects and their economic viability.\n
            File Type: xlxs\n\n
            Summary: Excel file presenting a detailed financial analysis of clean energy projects, evaluating their economic feasibility."""},
            {"role": "user", "content": f"Text:\n\n{summary_text_limit}\nFile Type:{doc_type}\n\nSummary:\n-"}]
        response = openai.ChatCompletion.create(
            engine='chatgpt',
            messages=messages,
            temperature=0,
            max_tokens=250,
            request_timeout=20
        )
        text = response.choices[0].message["content"]
        if text:
            return text.strip()
        return ""   
    except Exception as error:
        print(f"Failed->summary_generation. Error {error}")

def save_img(slide_list,slide_size,filename,container=False):
    image_url,thumb_url = [],[]
    for slide in slide_list:
        #save slide as an image to blob
        bmp = slide.get_thumbnail(1, 1)
        image_bytes = BytesIO()
        bmp.save(image_bytes, drawing.imaging.ImageFormat.png)
        image_bytes.seek(0)
        image_url.append(upload_image_thumbnail(f'{filename}/{str(slide.slide_number)}.png','image',image_bytes,container))
        if container:
            thumb_url.append(None)
            continue
        #save thumbnail version of the image to blob
        desiredX = slide_size.size.width/1.5
        desiredY = slide_size.size.height/1.5
        ScaleX = (1.0 / slide_size.size.width) * desiredX
        ScaleY = (1.0 / slide_size.size.height) * desiredY
        bmp = slide.get_thumbnail(ScaleX, ScaleY)
        thumbnail_bytes = BytesIO()
        bmp.save(thumbnail_bytes, drawing.imaging.ImageFormat.png)
        thumbnail_bytes.seek(0)
        thumb_url.append(upload_image_thumbnail(f'{filename}/{str(slide.slide_number)}.png','thumbnail',thumbnail_bytes,container))
    return (image_url,thumb_url)

def divide_range_into_parts(number):
    difference = 15
    ranges = []
    start = 0
    while start < number:
        end = min(start + difference, number)
        ranges.append((start, end))
        start = end
    return ranges

def ppt2img_thumb(path,container=False):
    status,file=download_blob(path)
    if status:
        filename = path.split("/")[-1].split(".")[0]
        presentation = slides.Presentation(BytesIO(file['file']))
        image_urls = []
        thumb_urls = []
        pres_size = presentation.slide_size
        call_split_into_images_with_args = lambda range_: save_img(presentation.slides[range_[0]:range_[1]],pres_size,filename,container)
        ranges = [(start, end) for start, end in divide_range_into_parts(presentation.slides.length) if start <= end]
        print("Preview generation started")
        with concurrent.futures.ThreadPoolExecutor(max_workers=15) as executor:
            for res in  executor.map(call_split_into_images_with_args,ranges):
                image_urls.extend(res[0])
                thumb_urls.extend(res[1])
        print("Preview generation started")
        
        return image_urls,thumb_urls
    return [],[]

def save_img_pdf(pdf_list,range,filename,container=False):
    image_urls,thumb_urls = [],[]
    for i,page in zip(range,pdf_list):
        image_bytes = BytesIO()
        thumbnail_bytes = BytesIO()
        if page.height-page.width < 300:
            new_image = Image.new('RGB', (1280, 1308), (255, 255, 255))
            page.thumbnail((1280,1308))
            padding = new_image.height-page.height
            top_padding =  padding// 2
            new_image.paste(page, (0, top_padding))

            new_image.save(image_bytes, 'PNG')
            new_image.thumbnail((354,500))  
            new_image.save(thumbnail_bytes, 'PNG')
        else:
            page.thumbnail((1280,1308))
            page.save(image_bytes, 'PNG')
            page.thumbnail((354,500))  
            page.save(thumbnail_bytes, 'PNG')
        image_bytes.seek(0)
        image_urls.append(upload_image_thumbnail(f'{filename}/{i}.png','image',image_bytes,container))
        if container:
            thumb_urls.append(None)
            continue
        thumbnail_bytes.seek(0)
        thumb_urls.append(upload_image_thumbnail(f'{filename}/{i}.png','thumbnail',thumbnail_bytes,container))
    return (image_urls,thumb_urls)


def pdf2img_thumb(path,container=False):
    status,file=download_blob(path)
    if status:
        filename = path.split("/")[-1].split(".")[0]
        pdf = convert_from_bytes(file['file'])
        image_urls = []
        thumb_urls = []
        call_split_into_images_with_args = lambda range_: save_img_pdf(pdf[range_[0]:range_[1]],range(range_[0],range_[1]),filename,container)
        ranges = [(start, end) for start, end in divide_range_into_parts(len(pdf)) if start <= end]
        print("Preview generation started")
        with concurrent.futures.ThreadPoolExecutor(max_workers=15) as executor:
            for res in  executor.map(call_split_into_images_with_args,ranges):
                image_urls.extend(res[0])
                thumb_urls.extend(res[1])
                print(thumb_urls)
        print("Preview generation ended")
        
        return image_urls,thumb_urls
    return [],[]

def check_recursively_for_imagetxt(shape):
    try:
        images = []
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                sub_images = check_recursively_for_imagetxt(s)
                if sub_images:
                    images.extend(sub_images)
        return images if images else [None]
    except Exception as error:
        print(f'Failed ppt_methods->check_recursively_for_imagetxt. Error:{error}')
        return [None]

@retry(Exception, tries=5, delay=1)
def find_title(prompt):
    try:
        rindex = random.randint(0,1)
        openai.api_base = openai_api_base[rindex]
        openai.api_key = openai_key[rindex]
        messages = [{"role": "user", "content": f"Generate a short title with at most 10 words on the below text delimited by ```. Here is the text:\n\n```{prompt}```"}]
        response = openai.ChatCompletion.create(
            engine="chatgpt",
            messages=messages,
            temperature=0,
            request_timeout=10
        )
        text = response.choices[0].message["content"]
        if text:
            return text.replace('"', '')
        return ""
    except RequestException:
        print("Request Exception")
        return ""
        
def check_recursively_for_text(shape):
    try:
        result = ("", "")
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP :
            for s in shape.shapes:
                sub_result = check_recursively_for_text(s)
                result = (result[0] + sub_result[0], result[1] + sub_result[1])
        else:
            shape_name = None
            if hasattr(shape, "text"):
                text = shape.text
                shape_name = shape.name.split()[0].replace(":","")
                if len(text)>0:
                    if shape_name == 'Title':
                        result = (result[0] + ". ".join(txt for txt in text.splitlines() if len(txt)!=0), result[1])
                    else:
                        result = (result[0], result[1] + ". ".join(txt  for txt in text.splitlines() if len(txt)!=0))
            
            if shape.has_table:
                text = ""
                table = shape.table
                for r in range(0,len(table.rows)):
                    for c in range(2,len(table.columns)):
                        cell_value = (table.cell(r,c)).text_frame.text
                        if cell_value:
                            text+=cell_value
                result = (result[0], result[1] + ". ".join(txt.strip()  for txt in text.splitlines() if len(txt)!=0))
        return result
    except Exception as error:
        print(f'Failed ppt_methods->check_recursively_for_text. Error:{error}')
        return ("", "")
    
def get_presentation_metadata(ppt):
    try:
        print("File Metadata extraction initiated.")
        presentation_meta = {}
        core_props = ppt.core_properties
        
        presentation_meta['author'] = core_props.author
        created_on = core_props.created
        presentation_meta['created_on'] = created_on.strftime('%Y-%m-%d %H:%M:%S') if created_on else ''
        presentation_meta['last_modified_by'] = core_props.last_modified_by
        modified_on = core_props.modified
        presentation_meta['modified_on'] = modified_on.strftime('%Y-%m-%d %H:%M:%S') if modified_on else ''
        presentation_meta['subject'] = core_props.subject
        presentation_meta['version'] = core_props.version
        if presentation_meta:
            return True, presentation_meta
        return False, "Failed to get ppt Metadata."
    except Exception as error:
        print(f'Falied ppt_methods->get_presentation_metadata. Error:{error}')

def get_slide_data(num,slide,workspace):
    slides_meta={}
    try:
        print(f"started processing for slide {num}")
        title = []
        content = []
        image_bytes = []
        images = []
        slides_meta['chunk_id'] = slide.slide_id
        slides_meta['chunk_number'] = num+1
        slides_meta["chunk_notes"] = ''
        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as excecutor:
            for res in excecutor.map(check_recursively_for_text,slide.shapes):
                title.append(res[0]),content.append(res[1])

        image_bytes = [image for shape in slide.shapes for image in check_recursively_for_imagetxt(shape) if image]


        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as excecutor:
            for res in excecutor.map(get_PPTimage_text,image_bytes):
                images.append(res)
                # time.sleep(1)
        
        title = [txt.strip() for txt in title if len(txt.strip())!=0]
        content = [txt.strip() for txt in content if len(txt.strip())!=0]
        images = [txt.strip() for txt in images if txt]
        slides_meta["chunk_title"],slides_meta["chunk_content"], slides_meta['image_text'] = ". ".join(title),". ".join(content), " ".join(images)
        
        if slide.has_notes_slide:
            slides_meta["chunk_notes"] = ". ".join(slide.notes_slide.notes_text_frame.text.splitlines())
        
        slides_meta['chunk_raw_text'] = " ".join([slides_meta["chunk_title"],slides_meta["chunk_content"],slides_meta['image_text'],slides_meta["chunk_notes"]])
        slides_meta['chunk_full_text'] = re.sub(r'([^a-zA-Z0-9.? ]+?)','',slides_meta['chunk_raw_text'])
        slides_meta['chunk_full_text'] = re.sub(r'\s+', ' ',slides_meta['chunk_full_text'])
        slides_meta['chunk_title'] = ""
        slides_meta['chunk_title'] = find_title(slides_meta['chunk_full_text'])
        slides_meta['chunk_title'] = slides_meta['chunk_title'] if slides_meta['chunk_title'] else ""

        dict_items = extract_ner(slides_meta['chunk_full_text'],workspace) if slides_meta['chunk_full_text'].strip() else {}
        slides_meta['entity'] = {k:list(v) for k,v in dict_items.items()}


        print(f"Finished processing for slide {num}")
        return slides_meta
    except Exception as error:
        print(f'Falied ppt_methods->slide_data. Error:{error}')
        return slides_meta

def check_file(path,workspace):
    try:
        db = Database('DocumentCollection',workspace)
        status,record = db.fetch_one_record({"file_path":path},{'_id':1,'edit_flag':1})
        if status:
            if record['edit_flag'] == 0:
                return delete_data(record['_id'],workspace)
            return False,"File already exists."
        
        client_status, blob_client = blob_connection()
        container = path.split("/")[3]
        print(container)
        blob_name = str(path.split(f"{container}/")[-1])
        message = "failed to create client."
        
        if client_status:
            print("Blob client successfully created")
            blob = blob_client.get_blob_client(
                container, blob_name)
            if not blob.exists():
                return True,None
        return delete_blob(path)
    except Exception as error:
        print(f'Failed -> check_file. Error:{error}')
        return False,message

def get_id(path,display_name,workspace):
    try:
        status,exists = check_file(path,workspace)
        if status:
            db = Database('DocumentCollection',workspace)
            presentation_meta = {}
            presentation_meta['file_path'] = path
            presentation_meta['file_name'] = path.split("/")[-1]
            presentation_meta['display_name'] = display_name.rsplit(".",1)[0]
            presentation_meta['file_source'] = 'upload'
            presentation_meta['edit_flag'] = 0

            status, id = db.insert_single_record(presentation_meta)
            message = id if not status else ""
            if status:
                return True, id
            return False, message
        return False, exists
    except Exception as error:
        print(f'Falied utils->get_id. Error:{error}')
        return False, "Failed to create record."

def preproc_list(items):
    unique_items = []
    for item in items:
        words = item.split()
        unique_words = [words[0]]
        for i in range(1, len(words)):
            if words[i] != words[i-1]:
                unique_words.append(words[i])
        add = {'atassigned','purchases','touchpoint','touchpoints','new','total','execution','com','media','consumer'}
        unique_item = " ".join([word for word in unique_words if word.lower() not in set(stopwords.words('english')).union(add) and word.strip()])
        if unique_item not in unique_items:
            unique_items.append(unique_item)
    return [words for words in unique_items if words]

def remove_string_subsets(original_list):
    new_list = []
    for s in original_list:
        if not any(s.lower() in p.lower() for p in new_list):
            new_list = [p for p in new_list if not p.lower() in s.lower()] + [s]
        elif s.lower() in [p.lower() for p in new_list]:
            continue
    return new_list

def summary_deduplication(text):
    sentences = text.split('\n-')
    filtered_sentences = []
    for sentence in sentences:
        is_duplicate = False
        for filtered_sentence in filtered_sentences:
            similarity_ratio = fuzz.token_set_ratio(sentence, filtered_sentence)
            if similarity_ratio >= 80:
                is_duplicate = True
                break
        if not is_duplicate:
            filtered_sentences.append(sentence)
    filtered_text = "-  "+('\n- '.join(filtered_sentences))
    return filtered_text

def get_ppt_data(blob_path,workspace):
    try:
        print("Process Started.")
        print("Started data extraction from PPT file.")
        status,file_data = download_blob(blob_path)
        message = file_data if not status else ""
        if status:
            ppt = Presentation(BytesIO(file_data['file']))
            status , presentation_meta = get_presentation_metadata(ppt)
            presentation_meta['uploaded_on']=file_data['created_on']
            message = presentation_meta if not status else ""
            if status:
                presentation_meta['file_type'] = 'ppt'
                presentation_meta['file_source'] = 'upload'
                presentation_meta['content'] = []
                presentation_meta['file_full_text'] = ''
                keywords = ppt.core_properties.keywords
                presentation_meta['tags'] = keywords.split(detect(keywords)) if len(keywords)!=0 else []
                print("File Metadata extracted.")

                slide_meta = []
                with concurrent.futures.ThreadPoolExecutor(15) as excecutor:
                    for res in  excecutor.map(get_slide_data,range(len(ppt.slides)),ppt.slides,repeat(workspace)):
                        slide_meta.append(res) 
                print('Slides data extraction compeleted.')
                # print(slide_meta)
                presentation_meta['file_title'] = ''
                if len(slide_meta)!=0:
                    if len(slide_meta[0]['chunk_title']) != 0:
                        if len(slide_meta[0]['chunk_title']) !=0:
                            presentation_meta['file_title'] = " ".join(txt.strip().capitalize() for txt in " ".join(slide_meta[0]['chunk_title'].split(".")).split())
                    else:
                        if len(slide_meta[0]['chunk_content']) != 0:
                            if len(slide_meta[0]['chunk_content']) !=0:
                                presentation_meta['file_title'] = " ".join(txt.strip().capitalize() for txt in " ".join(slide_meta[0]['chunk_content'].split(".")[:2]).split())
                for slide in slide_meta:
                    if not slide['chunk_title']:
                        slide['chunk_title'] = presentation_meta['file_title']
                presentation_meta['content'] = slide_meta
                presentation_meta['file_raw_text'] = " ".join(slides['chunk_raw_text'] for slides in slide_meta)
                presentation_meta['file_full_text'] = " ".join(slides['chunk_full_text'] for slides in slide_meta)
                
                presentation_meta['file_entities'] = {}
                temp_dict = [slides['entity'] for slides in slide_meta]
                for dictionary in temp_dict:
                    for key,value in dictionary.items():
                        if key not in presentation_meta['file_entities'].keys():
                            presentation_meta['file_entities'][key] = set()
                        value = [" ".join(val.split()).title() for val in value]
                        presentation_meta['file_entities'][key].update(value)
                presentation_meta['file_entities'] = [{"word":k.capitalize(),"label":res} for k,v in presentation_meta['file_entities'].items() if (res:=remove_string_subsets(preproc_list(list(v))))]
                
                presentation_meta['file_dictionary'] = dict()

                presentation_meta['number_of_chunks'] = len(slide_meta)
                presentation_meta['word_count'] = len(presentation_meta['file_full_text'].split())
                print("PPT file data extraction completed successfully.")
                
                if status:
                    print("Process Completed.")
                    print("PPT data has been updated.")
                    return True, presentation_meta
                print(message)
                return False, message
        print(f'Failed to read PPT data. Reason: {message}')
        return False, "Failed to read PPT data."
    except Exception as error:
        print(f'Failed ppt_methods->get_ppt_data. Error:{error}.')

def get_pdf_metadata(pdf):
    print("File Metadata extraction initiated.")
    pdf_metadata = {}
    meta = pdf.metadata
    
    pdf_metadata['author'] = meta.author
    created_on = meta.creation_date
    pdf_metadata['created_on'] = created_on.strftime('%Y-%m-%d %H:%M:%S') if created_on else ''
    created_on = meta.modification_date
    pdf_metadata['modified_on'] = created_on.strftime('%Y-%m-%d %H:%M:%S') if created_on else ''
    return True,pdf_metadata

def get_PDFimage_text(image):
    try:
        img = image.data
        text = ''
        # text = image_to_text(img)
        status,text = ocr_img2text(img)
        if status:
            if text:
                return text
        return ''
    except Exception as error:
        print(f'Failed ppt_methods->get_PDFimage_text. Error:{error}')
        return ''
    
def get_page_data(num,page,workspace):
    pages_meta={}
    try:
        print(f"started processing for page {num}")
        
        images = []
        pages_meta['chunk_number'] = num+1
        pages_meta['chunk_raw_text'] =''
        pages_meta['chunk_full_text'] =''
        pages_meta['chunk_title'] = ""
        pages_meta['entity'] = {}
        text = page.extract_text()
        page_text = text
        
        
        try:
            images_list = page.images
            if images_list:
            
                with concurrent.futures.ThreadPoolExecutor(max_workers=5) as excecutor:
                    for res in excecutor.map(get_PDFimage_text,images_list):
                        images.append(res)
        
            images = [re.sub(r'\s+', ' ', txt.strip()) for txt in images if txt]
        except:
            pass
  
        pages_meta['image_text'] = "\n".join(images)
        pages_meta['chunk_raw_text'] = re.sub(r'\s+', ' ', page_text.strip()) if page_text else ""
        pages_meta['chunk_raw_text'] = "\n".join([pages_meta['chunk_raw_text'],pages_meta['image_text']])
        
        pages_meta['chunk_full_text'] = re.sub(r'([^a-zA-Z0-9.? ]+?)','',pages_meta['chunk_raw_text'])
        pages_meta['chunk_full_text'] = re.sub(r'\s+', ' ',pages_meta['chunk_full_text'])

        pages_meta['chunk_title'] = find_title(pages_meta['chunk_full_text']) if pages_meta['chunk_full_text'] else ""
        pages_meta['chunk_title'] = pages_meta['chunk_title'] if pages_meta['chunk_title'] else ""

        dict_items = extract_ner(pages_meta['chunk_full_text'],workspace) if pages_meta['chunk_full_text'].strip() else {}
        pages_meta['entity'] = {k:list(v) for k,v in dict_items.items()}

        print(f"Finished processing for page {num}")
        return pages_meta
    except Exception as error:
        print(f'Falied ppt_methods->get_page_data page: {num}. Error:{error}')
        return pages_meta

def get_summary_tags(id,workspace):
    try:
        print("Summary and tag generation Started.")
        db = Database('DocumentCollection',workspace)
        print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Retrieving the record to be processed.")
        status,result = db.fetch_one_record({"_id" : id},{'file_type':1,'file_full_text':1,'file_raw_text':1})
        message = result if not status else ""
        if status:
            presentation_meta = {}
            presentation_meta['file_summary'] = ''
            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Slides Data and Metadata extracted.")
            summary = summary_generation(result['file_full_text'])
            summary = summary_deduplication(summary)
            presentation_meta['file_summary'] = summary
            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Summary Generated Successfully.")
            presentation_meta['doc_summary'] = doc_summary_generation(result['file_raw_text'],result['file_type'])
            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Document summary Generated Successfully.")
            presentation_meta['tags'],presentation_meta['word_cloud'] = get_keywords(result['file_full_text'])
            print(f"{datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')}: Word cloud Generated Successfully.")
            status,message = db.update_one_record(id,presentation_meta)
            if status:
                print("Process Completed.")
                print("PPT data has been updated.")
                return True, None
            print(message)
            return False, message
        print(f'Failed to generate summary and word cloud. Reason: {message}')
        return False, "Failed to generate summary and word cloud."
    except Exception as error:
        print(f'Failed ppt_methods->get_summary_tags. Error:{error}.')


def get_pdf_data(blob_path,workspace):
    try:
        print("Process Started.")
        print("Started data extraction from PPT file.")
        print("Retrieving the record to be processed.")
        status,file_data = download_blob(blob_path)

        message = file_data if not status else ""
        if status:
            pdf = PyPDF2.PdfReader(BytesIO(file_data['file']))
            status , pdf_meta = get_pdf_metadata(pdf)
            pdf_meta['uploaded_on']=file_data['created_on']
            message = pdf_meta if not status else ""
            if status:
                pdf_meta['file_type'] = 'pdf'
                pdf_meta['file_source'] = 'upload'
                pdf_meta['content'] = []
                pdf_meta['file_full_text'] = ''
                print("File Metadata extracted.")

                page_meta = []
                with concurrent.futures.ThreadPoolExecutor(15) as excecutor:
                    for res in  excecutor.map(get_page_data,range(len(pdf.pages)),pdf.pages,repeat(workspace)):
                        page_meta.append(res) 
                print('Pages data extraction compeleted.')
                # import json
                # with open("ner_output.txt", "w") as f:
                #     json.dump(ner_couts,f)
                pdf_meta['file_title'] = ''
                if page_meta:
                    pdf_meta['file_title'] = find_title(page_meta[0]['chunk_title'])
                
                for page in page_meta:
                    if not page['chunk_title']:
                        page['chunk_title'] = pdf_meta['file_title'] if pdf_meta['file_title'] else ""
                pdf_meta['content'] = page_meta
                pdf_meta['file_raw_text'] = " ".join(pages['chunk_raw_text'] for pages in page_meta)
                pdf_meta['file_full_text'] = " ".join(pages['chunk_full_text'] for pages in page_meta)
                
                pdf_meta['file_entities'] = {}
                temp_dict = [pages['entity'] for pages in page_meta]
                for dictionary in temp_dict:
                    for key,value in dictionary.items():
                        if key not in pdf_meta['file_entities'].keys():
                            pdf_meta['file_entities'][key] = set()
                        value = [" ".join(val.split()).title() for val in value]
                        pdf_meta['file_entities'][key].update(value)
                pdf_meta['file_entities'] = [{"word":k.capitalize(),"label":res} for k,v in pdf_meta['file_entities'].items() if (res:=remove_string_subsets(preproc_list(list(v))))]
                
                pdf_meta['file_dictionary'] = dict()

                pdf_meta['number_of_chunks'] = len(page_meta)
                pdf_meta['word_count'] = len(pdf_meta['file_full_text'])
                print("PDF file data extraction completed successfully.")
                
                if status:
                    print("Process Completed.")
                    print("PDF data has been updated.")
                    return True, pdf_meta
                print(message)
                return False, message
        print(f'Failed to read PDF data. Reason: {message}')
        return False, "Failed to read PDF data."
    except Exception as error:
        print(f'Failed pdf_methods->get_pdf_data. Error:{error}.')

def get_data(id,workspace):
    db = Database('DocumentCollection',workspace)
    print("Retrieving the record to be processed.")
    status,result = db.fetch_one_record({"_id" : id})
    message = result if not status else "" 
    if status:
        ext = result['file_path'].split('.')[-1]
        if ext in ['ppt','pptx']:
            status,presentation_meta = get_ppt_data(result['file_path'],workspace)
            if status:
                status,message = db.update_one_record(id,presentation_meta)
                if status:
                    return True, id
        elif ext in ['pdf']:
            status,presentation_meta = get_pdf_data(result['file_path'],workspace)
            if status:
                status,message = db.update_one_record(id,presentation_meta)
                if status:
                    return True, id
    return False, message

def split_ppt(id,workspace):
    try:
        print("Preview started")
        presentation_meta = {}
        db = Database('DocumentCollection',workspace)
        print("Retrieving the record to be processed.")
        status,result = db.fetch_one_record({"_id" : id})
        message = result if not status else ""
        if status:
            ext = result['file_path'].split(".")[-1]
            if ext in ['ppt','pptx']:

                images,thumbnails = ppt2img_thumb(result['file_path'])
            elif ext == 'pdf':
                images,thumbnails = pdf2img_thumb(result['file_path'])
            presentation_meta['file_preview'] = [{"chunk_number":num+1,"thumbnail_path":thumb,"image_path":img} for num,(thumb,img) in enumerate(zip(thumbnails,images))]
            status,message = db.update_one_record(id,presentation_meta)
            if status:
                print("Preview finished.")
                return True, message
        return False, message
    except Exception as error:
        print(f'Failed utils->split_ppt. Error:{error}.')
        return False, "Error while splitting ppt."

@cached(cache = TTLCache(maxsize=10, ttl=30), lock=Lock())
def get_categorywise_ppt(workspace,filter = None,user_id = None):
    try:
        db = Database('UserActivity',workspace)
        exclude_ids = []
        status,record = db.fetch_one_record({'user_id':user_id},{'listing_filter':1})
        if status:
            if 'listing_filter' in record:
                exclude_ids = record['listing_filter']
        print(exclude_ids)
        db = Database('DocumentCollection',workspace)
        query = [
                {'$match':{'edit_flag':1,"_id": {"$nin": exclude_ids}}},
                {'$project': {'_id': 1,'display_name': 1,'file_title': 1,'file_name':1, 'file_type':1, 'file_path':1, 'file_preview': 1,'file_summary':1,'word_cloud':1,'category': {'$objectToArray': '$content_categories'}, 'modified_on':1}},
                {'$unwind': {"path": '$category',"preserveNullAndEmptyArrays": True}},
                {'$project': {'_id': 1,'display_name': 1,'file_title': 1,'file_name':1, 'file_type':1, 'file_path':1, 'file_preview': 1,'file_summary':1,'word_cloud':1,'category': {'$ifNull': [ '$category.k', 'other' ] },'count': {'$ifNull': [ '$category.v', 0 ] }, 'modified_on':1}}
            ]
        status,results = db.fetch_aggregate(query)
        
        if status:
            records = pd.DataFrame(results)
            grouped = records.groupby('_id')
            ids_to_update = grouped.filter(lambda x: (x['count'] >= 5).any() == False)['_id']
            records.loc[records['_id'].isin(ids_to_update), ['category','count']] = ['Other',0]

            records = records[(records['count'] >= 5) | (records['category'] == 'Other')]

            records = records.drop_duplicates(subset=['_id', 'category'])
            if user_id:
                db1 = Database('UserActivity',workspace)
                query = [
                    {'$match': {'user_id': user_id}}, 
                    {'$unwind': {'path': '$file_history'}}, 
                    {'$lookup': {'from': 'DocumentCollection','localField': 'file_history.file_id','foreignField': '_id', 'as': 'file'}}, 
                    {'$unwind': {'path': '$file'}}, 
                    {"$addFields": {"category": "Previously Viewed"}},  # add a new field with a given value
                    {'$project': {'_id': '$file._id','display_name': '$file.display_name','file_title': '$file.file_title','file_name': '$file.file_name', 'file_type':'$file.file_type', 'file_path': '$file.file_path', 'file_preview': '$file.file_preview','file_summary': '$file.file_summary','word_cloud': '$file.word_cloud','category': 1, 'modified_on':'$file.modified_on'}}
                ]
                status,result = db1.fetch_aggregate(query)
                if status:
                    result = pd.DataFrame(result)
                    result['count'] = result.index+1
                    records = pd.concat([result,records])
            query = [
                {"$match": {"edit_flag": 1, "click_count": {"$exists": True},"_id": {"$nin": exclude_ids}}},  # filter by edit_flag = 1
                {"$project": {'_id': 1,'display_name': 1,'file_title': 1,'file_name':1, 'file_type':1, 'file_path':1, 'file_preview': 1,'file_summary':1,'word_cloud':1, 'click_count':1, 'modified_on':1}},  # project specified fields
                {"$sort": {"click_count": -1}},  # sort by click_count in descending order
                {"$addFields": {"category": "Trending Content","count": "$click_count"}},  # add a new field with a given value
                {"$match": {"count": {"$gt": 10}}},  # filter for sorted_click_count > 10
                {"$limit": 50},  # limit to top 50 records
                {"$project": {'_id': 1,'display_name': 1,'file_title': 1,'file_name':1, 'file_type':1, 'file_path':1, 'file_preview': 1,'file_summary':1,'word_cloud':1, 'category':1, 'count':1, 'modified_on':1}},
            ]
            status,result = db.fetch_aggregate(query)
            if status:
                records = pd.concat([pd.DataFrame(result),records])
            if filter:
                records = records[records['category']==filter].drop(['category'],axis=1).sort_values(['count', 'modified_on'],ascending=[False,False]).reset_index(drop=True)
                return True, records
            return True,records
        else:
            return False,results
    except Exception as error:
        print(f"Failed->get_categorywise_ppt. Error:{error}")

def get_home_page_ppts(user_id,workspace):
    try:
        status,records = get_categorywise_ppt(user_id = user_id,workspace=workspace)
        if status:
            data = []
            category_list = records['category'].unique()
            category_list = list(category_list)
            category_list = sorted(category_list)
            if "Other" in category_list: 
                category_list.remove('Other')
                category_list.append('Other')
            if "Trending Content" in category_list:
                category_list.remove("Trending Content")
                category_list = ['Trending Content'] + category_list
            if "Previously Viewed" in category_list:
                category_list.remove("Previously Viewed")
                category_list = ['Previously Viewed'] + category_list
            for category in category_list:
                cats = {}
                cats['category'] = category
                files = records[records['category']==category].drop(['category'],axis=1).sort_values(['count', 'modified_on'],ascending=[False,False]).head(20)
                cats['files'] = files.to_dict('records')
                data.append(cats)
            return True,data
        else:
            return False,records
    except Exception as error:
        print(f"Failed->get_home_page_ppts. Error:{error}")

def get_see_more_ppts(category,page_num,user_id,workspace):
    try:
        status,records = get_categorywise_ppt(filter=category,user_id = user_id,workspace=workspace)
        if status:
            data = []
            cats = {}
            cats['category'] = category
            files = records
            start_index = (int(page_num) - 1) * 10
            end_index = start_index + 10
            page = files.iloc[start_index:end_index]
            cats['files'] = page.to_dict('records')
            data.append(cats)
            return True,data
        else:
            return False,records
    except Exception as error:
        print(f"Failed->get_see_more_ppts. Error:{error}")

def process_categories(text,workspace):
    status, breakouts = get_user_settings(col={'_id':0,'breakouts':1},workspace=workspace)
    if status:
        categories = {category['category']:category['tags'] for category in breakouts['breakouts']}
    categories_keywords = {}
    content_categories = {}
    words = Counter(text.lower().split())
    for category in categories:
        keywords = {word:words.get(word) for word in set(categories[category]).intersection(set(words.keys()))}
        categories_keywords[category] = [{"word":k, "weight":v} for k,v in keywords.items()]
        content_categories[category] = sum(keywords.values())
    return categories_keywords,content_categories


def get_categories(id,workspace):
    try:
        db = Database('DocumentCollection',workspace)
        print("Retrieving the record to be processed.")
        status,result = db.fetch_one_record({"_id" : id},{"file_full_text":1})
        message = result if not status else ""
        if status:
            ppt_categories = {}
            ppt_categories['categories_keywords'],ppt_categories['content_categories'] = process_categories(result['file_full_text'],workspace)
            status,message = db.update_one_record(id,ppt_categories)
            if status:
                print("PPT category data has been updated.")
                return True, None
        print(message)
        return False, message
    except Exception as error:
        print(f"Failed->get_categorywise_ppt. Error:{error}")

def increment_onclick(id,workspace,category = False):
    try:
        db = Database('DocumentCollection',workspace)
        if category:
            status,records = db.increment_one_record(id,optional_field=category)
        else:
            status,records = db.increment_one_record(id)
        if status:
            return True,records
        return False,records
    except Exception as error:
        print(f"Failed->increment_onclick. Error:{error}")

def semantic_search(search_query,index,workspace) -> str:
    if index==search_master_table:
        columns = ["oid", "content", "file_name", "file_type","title", "primary_id"]
    else:
        columns=['oid', 'content', 'chunk_number', 'title', 'file_type', 'display_name', 'primary_id']
    embedding_url = f"{DATABRICKS_URL}/serving-endpoints/databricks-bge-large-en/invocations"
    payload = json.dumps({"input": [search_query]})
    headers = {'Content-Type': 'application/json','Authorization': f"Bearer {PAT_TOKEN}"}
    response = requests.request("POST", embedding_url, headers=headers, data=payload)
    print('Embedding response: ', response)
    data=response.json()['data'][0]['embedding']
    workspace = workspace.replace('-','_')
    search_url = f"{DATABRICKS_URL}/api/2.0/vector-search/endpoints/{workspace}/indexes/{CATALOG}.{workspace}.{index}_index/query"
    payload = json.dumps({
        "query_vector": data,
        "columns": columns,
        "num_results": 50
        })
    headers = {'Content-Type': 'application/json','Authorization': f"Bearer {PAT_TOKEN}"}
    response = requests.request("GET", search_url, headers=headers, data=payload)
    print('Search Response: ', response)
    return response.json()


def search_score(query, index_name, workspace):
    results = semantic_search(query,index_name,workspace)
    # scores = np.array([hit["@search.reranker_score"] if hit["@search.reranker_score"] else 0 for hit in results ])
    column_count=int(results['manifest']['column_count'])
    scores = np.array([hit[column_count-1] if hit[column_count-1] else 0 for hit in results['result']['data_array'] ])

    if not results:
        return pd.DataFrame()
    
    mean_score = np.mean(scores)
    std_score = np.std(scores)

    def sigmoid(x):
        return 1 / (1 + np.exp(-x))

    epsilon = 1e-8
    std_score += epsilon
    
    normalized_scores = (scores - mean_score) / std_score
    normalized_scores = sigmoid(normalized_scores)
   
    if index_name==search_page_table:
        normalized_results = [
            {
                "_id": hit[0],
                "Similarity_score": normalized_score,
                "chunk_number": hit[2],
                "chunk_title": hit[3]
            }
            # for hit, normalized_score in zip(results, normalized_scores)
            for hit, normalized_score in zip(results['result']['data_array'], normalized_scores)
        ]
    else:
        normalized_results = [
            {
                "_id": hit[0],
                "Similarity_score": normalized_score,
                # "chunk_number": hit[2] if index_name == search_page_table else None,
                "file_title": hit[4]
            }
            # for hit, normalized_score in zip(results, normalized_scores)
            for hit, normalized_score in zip(results['result']['data_array'], normalized_scores)
        ]

    return pd.DataFrame(normalized_results)

@cached(cache = TTLCache(maxsize=10, ttl=30), lock=Lock())
def get_ppt_search_result(search_query,user_id,workspace):
    try:
        db = Database('UserActivity',workspace)
        exclude_ids = []
        status,record = db.fetch_one_record({'user_id':user_id},{'listing_filter':1})
        if status:
            if 'listing_filter' in record:
                exclude_ids = record['listing_filter']
        # index = search_index.format(workspace)
        index = search_master_table
        db = Database('DocumentCollection',workspace)
        results  = search_score(search_query,index,workspace)
        if len(results)!=0:
            status,data = db.fetch_all_records({"edit_flag":1,"_id": {"$nin": exclude_ids}},{'_id': 1,'display_name': 1,'file_title': 1,'file_name':1, 'file_type':1, 'file_path':1, 'file_preview': 1,'file_summary':1,'word_cloud':1, 'modified_on':1})
            if not status:
                return False,"No data found"

            data = pd.DataFrame(data)
            results['_id'] = results['_id'].apply(lambda x: ObjectId(x))
            results['match_score'] = results['Similarity_score'].apply(lambda x:int(x*100))
            # results = results[['_id','match_score']]
            results = results[['_id','match_score']]
            final_results = pd.merge(results,data,how='inner',on=['_id'])
            if len(final_results)!=0:
                final_results = final_results.sort_values(by=['match_score', 'modified_on'], ascending=[False, False])
                final_results = final_results[final_results['match_score']>20]
                return True, final_results
        return True,pd.DataFrame()
    except Exception as error:
        print(f"Failed->get_ppt_search_result. Error:{error}")

@cached(cache = TTLCache(maxsize=10, ttl=30), lock=Lock())
def get_slide_search_result(search_query,user_id,workspace):
    try:
        db = Database('UserActivity',workspace)
        exclude_ids = []
        status,record = db.fetch_one_record({'user_id':user_id},{'listing_filter':1})
        if status:
            if 'listing_filter' in record:
                exclude_ids = record['listing_filter']
        # index = chunk_index.format(workspace)
        index = search_page_table
        db = Database('DocumentCollection',workspace)
        results  = search_score(search_query,index,workspace)
        if len(results)!=0:
            # results['chunk_number'] = results['_id'].str.split('_').str[1]
            results['_id'] = results['_id'].str.split('_').str[0]
            results['chunk_number'] = results['chunk_number'].astype(int)
            query = [
                    {'$match': {'edit_flag': 1}},
                    {'$project': {'display_name':1,'modified_on': 1, 'file_type':1, 'combined': {'$zip': {'inputs': ['$file_preview']}}}}, 
                    {'$unwind': {'path': '$combined'}}, 
                    {'$project': {'_id': 1, 'display_name':1,'modified_on': 1, 'file_type':1, 'chunk_number': {'$arrayElemAt': ['$combined.chunk_number', 0]},'thumbnail_path': {'$arrayElemAt': ['$combined.thumbnail_path', 0]}, 'image_path': {'$arrayElemAt': ['$combined.image_path', 0]}, 'download_path': ''}}
                ]
            status,data = db.fetch_aggregate(query)
            if not status:
                return False,"No data found"

            data = pd.DataFrame(data)
            # data['chunk_number'] = data['chunk_number'].astype(int)
            data['selected'] = False
            results['_id'] = results['_id'].apply(lambda x: ObjectId(x))
            results['match_score'] = results['Similarity_score'].apply(lambda x:int(x*100))
            # results = results[['_id','match_score','chunk_number']]
            # final_results = pd.merge(results,data,how='inner',on=['_id','chunk_number'])
            results = results[['_id','match_score','chunk_number','chunk_title']]
            final_results = pd.merge(results,data,how='inner',on=['_id', 'chunk_number'])
            if len(final_results)!=0:
                final_results = final_results.drop_duplicates(['chunk_title','match_score'])
                final_results = final_results.sort_values(by=['match_score', 'modified_on'], ascending=[False, False])
                final_results = final_results[final_results['match_score']>20]
                return True, final_results
        return False,"Empty Data"
    except Exception as error:
        print(f"Failed->get_slide_search_result. Error:{error}")

def add_search_history(search_term,user_id,workspace):
    try:
        db = Database('UserActivity',workspace)
        query = {'user_id':user_id}
        data = {'$pull':{"search_history":{"word":search_term.title()}}}
        status,records = db.append_data(query,data)
        data = {'$addToSet':{"search_history":{"word":search_term.title()}}}
        status,records = db.append_data(query,data)
        if status:
            return True,records
        return False,records
    except Exception as error:
        print(f"Failed->add_search_history. Error:{error}")

def get_search_results(query,page_num,user_id,workspace):
    try:
        status,records = get_ppt_search_result(query,user_id,workspace)
        #code to add query term in userActivity
        add_search_history(query,user_id,workspace)
        if status:
            data = []
            result = {}
            result['search_term'] = query
            files = records
            start_index = (int(page_num) - 1) * 10
            end_index = start_index + 10
            page = files.iloc[start_index:end_index]
            result['files'] = page.to_dict('records')
            result['totalRecord'] = len(files)
            data.append(result)
            return True,data
        return False,[]
    except Exception as error:
        print(f"Failed->get_search_results. Error:{error}")

def get_slide_search_results(workspace,query,page_num,user_id,filter = None):
    try:
        status,records = get_slide_search_result(query,user_id,workspace)
        add_search_history(query,user_id,workspace)
        if status:
            response = {}
            if filter:
                return True,records[records['display_name']==filter].drop(['display_name'],axis=1).sort_values(['match_score','modified_on'],ascending=[False,False]).reset_index(drop=True)
            sorted_df = records.sort_values(by=['_id', 'match_score'], ascending=[True, False])
            def set_priority(row):
                max_score = sorted_df.loc[sorted_df['_id'] == row['_id'], 'match_score'].max()
                group_size = sorted_df.loc[sorted_df['_id'] == row['_id']].shape[0]
                return max_score, group_size
            sorted_df['priority'] = sorted_df.apply(set_priority, axis=1)
            sorted_df = sorted_df.sort_values(by='priority', ascending=False)
            records = sorted_df.drop('priority', axis=1)
            ppts = records['display_name'].unique()
            
            start = (int(page_num)-1)*5
            end = min(len(ppts),start+5)
            ppts_page = ppts[start:end]
            data = []
            for ppt in ppts_page:
                subset = records[records['display_name']==ppt].drop(['display_name'],axis=1).sort_values(['match_score','modified_on'],ascending=[False,False]).reset_index(drop=True)
                result = {}
                result['modified_on'] = subset['modified_on'].unique()[0]

                if len(subset)!=0:
                    result['_id'] = subset['_id'].iloc[0]
                    result['display_name'] = ppt
                    result['file_type'] = subset['file_type'].iloc[0]
                    result['files'] = subset.to_dict('records')
                    result['pptRecord'] = len(subset)
                    data.append(result)
            response['search_term'] = query
            response['data'] = data
            response['totalRecord'] = len(records)
            response['loadingNext'] = False
            response['noMoreData'] = False if end < (len(ppts)) else True
            response['page_no'] = int(page_num)
            return True,response
        return False,records
    except Exception as error:
        print(f"Failed->get_slide_search_results. Error:{error}")

def date_category(row):
    return str(row['modified_on'].year)

def get_vintage_search_results(workspace,query,type,user_id,filter=None):
    try:
        if type == 'ppt_vintage':
            status,records = get_ppt_search_result(query,user_id,workspace)
            add_search_history(query,user_id,workspace)
        elif type == 'slide_vintage':
            status,records = get_slide_search_result(query,user_id,workspace)
        if status:
            if len(records)!=0:
                print("here")
                global current_date
                current_date = pd.Timestamp.now()
                records['modified_on'] = pd.to_datetime(records['modified_on'])
                records['category'] = records.apply(date_category, axis=1)
                records['date'] = records['modified_on'].dt.strftime('%Y-%m-%d')
                records['modified_on'] = records['modified_on'].dt.strftime('%Y-%m-%d %H:%M:%S')
                # if len(records)!=0:
                if filter:
                    records = records[records['category']==filter].drop(['category'],axis=1).sort_values(['modified_on','match_score'],ascending=[False,False]).reset_index(drop=True)
                    return True,records
                return True,records
            return False,"Empty Data"
        return False,"Empty Data"
    except Exception as error:
        print(f"Failed->get_vintage_search_results. Error:{error}")
        return False,pd.DataFrame()

def paginate_vintage(workspace,search_term,type,user_id,page_num=1,filter=None):
    try:
        status,record = get_vintage_search_results(workspace,search_term,type,user_id)
        print(status,record)
        if status:
            category_order = sorted(record['category'].unique(),reverse=True)
            if filter:
                category_order = [filter]
            response = {}
            data = []
            page_num = page_num if page_num else 1
            for category in category_order:
                subset = record[record['category']==category].sort_values(['date','match_score'],ascending=[False,False]).drop(['category','date'],axis=1).reset_index(drop=True)
                result = {}
                files = subset
                start_index = (int(page_num) - 1) * 5 if page_num else 0
                end_index = start_index + 5
                page = files.iloc[start_index:end_index]
                if len(page)!=0:
                    result['category'] = category
                    result['files'] = page.to_dict('records')
                    result['loadingNext'] = False
                    result['noMoreData'] = False if end_index < (len(files)) else True
                    result['page_no'] = page_num
                    data.append(result)
            response['data'] = data
            response['totalRecord'] = len(record)
            return True,response
        return False,record
    except Exception as error:
        print(f"Failed->paginate_vintage. Error:{error}")

def paginate_vintage_slides(workspace,search_term,type,user_id,page_num=1,filter=None,subfilter=None):
    try:
        status,record = get_vintage_search_results(workspace,search_term,type,user_id)
        print(status,record)
        if status:
            category_order = sorted(record['category'].unique(),reverse=True)
            response = {}
            response['search_term'] = search_term
            if filter:
                category_order = [filter]
            page_num = page_num if page_num else 1
            data = []
            for category in category_order:
                records = record[record['category']==category].drop(['category'],axis=1).reset_index(drop=True)
                sub_response = {}
                ppt_data = []
                if subfilter:
                    return True,records[records['_id']==subfilter].drop(['display_name'],axis=1).sort_values(by='match_score',ascending=False).reset_index(drop=True)
                sorted_df = records.sort_values(by=['_id', 'match_score'], ascending=[True, False])
                def set_priority(row):
                    max_score = sorted_df.loc[sorted_df['_id'] == row['_id'], 'match_score'].max()
                    group_size = sorted_df.loc[sorted_df['_id'] == row['_id']].shape[0]
                    return max_score, group_size
                sorted_df['priority'] = sorted_df.apply(set_priority, axis=1)
                sorted_df = sorted_df.sort_values(by=['date','priority'], ascending=[False,False])
                records = sorted_df.drop(['priority','date'], axis=1)
                ppts = records['display_name'].unique()
                start = (int(page_num)-1)*5
                end = min(len(ppts),start+5)
                ppts_page = ppts[start:end]
                for ppt in ppts_page:
                    subset = records[records['display_name']==ppt].drop(['display_name'],axis=1).reset_index(drop=True)
                    result = {}
                    if len(subset)!=0:
                        result['modified_on'] = subset['modified_on'].iloc[0]
                        result['_id'] = subset['_id'].iloc[0]
                        result['display_name'] = ppt
                        result['files'] = subset.to_dict('records')
                        result['pptRecord'] = len(subset)
                        ppt_data.append(result)
                sub_response['category'] = category
                sub_response['data'] = ppt_data
                sub_response['loadingNext'] = False
                sub_response['noMoreData'] = False if end < (len(ppts)) else True
                sub_response['page_no'] = int(page_num)
                data.append(sub_response)
            response['data'] = data
            response['totalRecord'] = len(record)
            return True,response
        return False,"Empty Data"
    except Exception as error:
        print(f"Failed->paginate_vintage_slides. Error:{error}")


def add_ppt_history(id,user_id,workspace):
    try:
        db = Database('UserActivity',workspace)
        query = {'user_id':user_id}
        data = {'$pull':{"file_history":{"file_id":id}}}
        status,records = db.append_data(query,data)
        data = {'$addToSet':{"file_history":{"file_id":id}}}
        status,records = db.append_data(query,data)
        if status:
            return True,records
        return False,records
    except Exception as error:
        print(f"Failed->add_ppt_history. Error:{error}")


def get_relavent_ppts(id, user_id,workspace, category = False, search = False,type = False):
    try:
        if search:
            if type=='ppt':
                status,records = get_ppt_search_result(search_query=search,user_id=user_id,workspace=workspace)
                increment_onclick(id,workspace=workspace)
                add_ppt_history(id,user_id,workspace)
            elif type=='ppt_vintage':
                status,records = get_vintage_search_results(workspace,search,type,user_id,category)
                print(status)
                increment_onclick(id,workspace=workspace)
                add_ppt_history(id,user_id,workspace)
        elif category:
           status,records = get_categorywise_ppt(filter=category,user_id=user_id,workspace=workspace)
           increment_onclick(id,workspace=workspace,category=category)
           add_ppt_history(id,user_id,workspace)
        if status:
            target_index = records.index[records['_id'] == id][0]
            start_index = max(0, target_index - (13 // 2))
            end_index = min(len(records), target_index + (13 // 2) + (13 % 2))
            num_preceding = target_index - start_index
            num_following = end_index - target_index - 1
            if num_preceding < (13 // 2):
                num_following += (13 // 2) - num_preceding
            if num_following < (13 // 2):
                num_preceding += (13 // 2) - num_following
            start_index = max(0, target_index - num_preceding)
            end_index = min(len(records), target_index + num_following + 1)
            files = records.iloc[start_index:end_index]
            files = files.drop(target_index).reset_index(drop=True) 
            return True,files.to_dict('records')
        else:
            return False,records
    except Exception as error:
        print(f"Failed->get_relavent_ppts. Error:{error}")

def get_relavent_slides(workspace,id, user_id, category = False, search = False,type = False, slide_number=False):
    try:
        if type=='slide':
            status,records = get_slide_search_results(workspace,search,None,user_id,filter=category)
            increment_onclick(id,workspace=workspace)
            add_ppt_history(id,user_id,workspace)
        elif type=='slide_vintage':
            status,records = paginate_vintage_slides(workspace,search,type,user_id,None,filter=category,subfilter=id)
        if status:
            print("id: ",id,"\nchunk_number :",slide_number)
            print(records)
            target_index = records.index[(records['_id'] == id)&(records['chunk_number']==slide_number)][0]
            start_index = max(0, target_index - (13 // 2))
            end_index = min(len(records), target_index + (13 // 2) + (13 % 2))
            num_preceding = target_index - start_index
            num_following = end_index - target_index - 1
            if num_preceding < (13 // 2):
                num_following += (13 // 2) - num_preceding
            if num_following < (13 // 2):
                num_preceding += (13 // 2) - num_following
            start_index = max(0, target_index - num_preceding)
            end_index = min(len(records), target_index + num_following + 1)
            files = records.iloc[start_index:end_index]
            files = files.drop(target_index).reset_index(drop=True) 
            return True,files.to_dict('records')
        return False,records
    except Exception as error:
        print(f"Failed->get_relavent_slides. Error:{error}")


def create_search_index(index_name):
    print(f"Ensuring search index {index_name} exists")
    if index_name not in searchindexclient.list_index_names():
        index = SearchIndex(
            name=index_name,
            fields=[
                SimpleField(name="id", type="Edm.String", key=True),
                SearchField(name="content", type="Edm.String", analyzer_name="en.microsoft",
                            hidden=False, searchable=True, filterable=False, sortable=False, facetable=False),    
                SearchField(name="file_name", type="Edm.String",analyzer_name="en.microsoft",
                            hidden=False, searchable=True, filterable=True, sortable=False, facetable=False)     
            ],
            semantic_settings=SemanticSettings(
                configurations=[SemanticConfiguration(
                    name='qa_semantic',
                    prioritized_fields=PrioritizedFields(
                        title_field=None, prioritized_content_fields=[SemanticField(field_name='content')]))])
            )
        print(f"Creating {index_name} search index")
        searchindexclient.create_index(index)
    else:
        print(f"Search index {index_name} already exists")


def get_search_recommendations(workspace):
    try:
        db = Database('DocumentCollection',workspace)
        query = [
            {'$match': {'edit_flag': 1}}, 
            {'$project': {'_id': 0,'word_cloud': 1,'file_entities':1,'file_title':1,'display_name':1}}
        ]
        status,results = db.fetch_aggregate(query)
        if status:
            merged = defaultdict(float)

            for items in results:
                for item in items['word_cloud']:
                    merged[item['word'].title()] = item['weight']
                for i in range(len(items['file_entities'])):
                    for label in items['file_entities'][i]['label']:
                        merged[label.title()]=1
                merged[items['file_title'].title()] = 1

            merged = [{'word': k, 'weight': v} for k, v in merged.items()]
            return True,merged
    except Exception as error:
        print(f"Failed:get_search_recommendations. Error: {error}")
        return False,results


def get_search_history(user_id,workspace):
    try:
        db = Database('UserActivity',workspace)
        filter = {"user_id":user_id}
        query = {'_id':0,"search_history":1}
        results = {
            'search_history':[],
            'search_recommendations':[]
            }
        status,record = db.fetch_one_record(filter,query)
        if status:
            results['search_history'] = record['search_history'][-1:-11:-1]

        status,record = get_search_recommendations(workspace)
        if status:
            results['search_recommendations'] = record
        return True,results
    except Exception as error:
        print(f"Failed:get_search_history. Error: {error}")
        return False,results

def get_presentation_data(id,search_term,user_id,workspace):
    try:
        status,data = get_ppt_search_result(search_term,user_id,workspace)
        if status:
            return True,data[data['_id']==id].to_dict('records')[0]
        return True,{}
    except Exception as error:
        print(f"Failed:get_presentation_data. Error: {error}")
        return False,'Failed to retrieve presentation data.'
    
def clean_up(workspace='smartsearch-exponentia-dev'):
    us_timezone = pytz.timezone('EST')
    current_time = datetime.now(us_timezone)
    db = Database('DocumentCollection',workspace)
    status,records = db.fetch_all_records({"edit_flag":0},{'_id':1})
    if status:
        if records:
            for record in records:
                delete_data(record['_id'],workspace)
    print(f"Job ran at: {current_time}")
    return

def categorize_all_records(workspace):

    db_central=Database("workspace","central")
    db=Database("DocumentCollection",workspace)
    job_id = None
    status,records=db_central.fetch_one_record({"workspace_name":workspace},{"_id":0,"upload_job_id":1})
    if status and "upload_job_id" in records:
        print('Categorizing records')
        job_id = records['upload_job_id']
    if job_id:
        workspaceurl=DATABRICKS_URL
        databricks_api_url = f"{workspaceurl}/api/2.1/jobs/run-now"
        token = PAT_TOKEN
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {token}"
        }
        payload = {
            "job_id": job_id,
            "job_parameters": {"type":"get_categories", "workspace":workspace}
        }
 
        logging.info("Sending request to Databricks...")
        response = requests.post(databricks_api_url, headers=headers, json=payload)
        if response.status_code == 200:
            try:
                notebook_response = response.json()
                run_id = notebook_response.get("run_id")
                url = f"{workspaceurl}/api/2.1/jobs/runs/get?run_id={run_id}"
                payload = ""
                headers = {"Authorization": f"Bearer {token}"}
        
                response = requests.request("GET", url, headers=headers, data=payload)
                if response.status_code==200:
                    notebook_response = response.json()
                    try:
                        task_id= list(filter(lambda x: x["task_key"] == "Fetch_all_categories", notebook_response["tasks"]))[0]['run_id']
                    except Exception as e:
                        print(e)
                    url = f"{workspaceurl}/api/2.1/jobs/runs/get-output?run_id={task_id}"
                    status = notebook_response["state"]["life_cycle_state"]
       
                    while status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                        response = requests.request("GET", url, headers=headers, data=payload)
                        if response.status_code==200:
                            notebook_response = response.json()
                            status = notebook_response['metadata']["state"]["life_cycle_state"]
 
                        else:
                            print(json.dumps({"notebook_response": "Error in fetching id"}))
                        if status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                            time.sleep(5)
                    data=notebook_response
                    job_status=notebook_response['metadata']["state"]["result_state"]
                    if job_status=='SUCCESS':
                        print("Successfully got all categories")
                    else:
                        print("Error In Getting Categories")
                else:
                    print('Unable to get categories')
                return
            except Exception as e:
                print(f"Error: categorize_all_records->{e}")
                return

def update_settings(id,suggestions,breakouts,ner,workspace):
    try:
        db = Database('UserConfiguration',workspace)
        settings = {
            'suggestions':suggestions,
            'breakouts':breakouts,
            'ner':ner}
        status,message = db.update_one_record(id,settings)
        if status:
            thread = threading.Thread(target=categorize_all_records(workspace))
            thread.start()
            return True,message
        return False,message
    except Exception as error:
        print(f"Failed:update_settings. Error: {error}")
        return False,"Failed to update settings."

def generate_cron_expression(frequency, start_time=None,day=None):
    cron_expression = None
    if frequency == "Hourly":
        timezone=pytz.timezone('Asia/Kolkata')
        hour=int((datetime.now(timezone)+timedelta(hours = 2)).strftime("%H"))
        minute=int(datetime.now(timezone).strftime("%M"))
        cron_expression = f"{hour} {minute} * * * ?"
    elif frequency == "Daily":
        hour,minute = start_time.split(":")
        cron_expression = f"0 {minute} {hour} * * ?"
    elif frequency == "Weekly":
        hour,minute = start_time.split(":")
        cron_expression = f"0 {minute} {hour} ? * {day.title()}"
    else:
        raise ValueError("Unsupported frequency")
    return cron_expression

def create_job(pipeline,workspace):
    pipeline_name = pipeline['name']
    site_id = pipeline['siteid']
    folder_name = pipeline['foldername']
    job_config = {
        "name": f"{workspace}_{pipeline_name}",
        "email_notifications": {
            "no_alert_for_skipped_runs": False},
        "webhook_notifications": {},
        "timeout_seconds": 0,
        "schedule": {
            "quartz_cron_expression": generate_cron_expression(pipeline['frequency'].title(),pipeline['time'],pipeline['day']),  # Every day at midnight
            "timezone_id": "Asia/Kolkata"},
        "max_concurrent_runs": 1,
        "tasks": [
            {
                "task_key": "bronze_silver",
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"},
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False},
                "webhook_notifications": {}
            },
            {
                "task_key": "gold_layer_chat",
                "depends_on": [{"task_key": "bronze_silver"}],
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"},
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False},
                "webhook_notifications": {}
            },
            {
                "task_key": "gold_layer_search",
                "depends_on": [{"task_key": "bronze_silver"}],
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"},
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False},
                "webhook_notifications": {}
            },
            {
                "task_key": "Chat_Master_Index",
                "depends_on": [{"task_key": "gold_layer_chat"},{"task_key": "gold_layer_search"}],
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"},
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False},
                "webhook_notifications": {}
            },
            {
                "task_key": "Chat_Repo_Index",
                "depends_on": [{"task_key": "gold_layer_chat"},{"task_key": "gold_layer_search"}],
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"},
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False
                },
                "webhook_notifications": {}
            },
            {
                "task_key": "Search_Master_Index",
                "depends_on": [{"task_key": "gold_layer_chat"},{"task_key": "gold_layer_search"}],
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"
                },
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False
                },
                "webhook_notifications": {}
            },
            {
                "task_key": "Search_Repo_Index",
                "depends_on": [{"task_key": "gold_layer_chat"},{"task_key": "gold_layer_search"}],
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"
                },
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False},
                "webhook_notifications": {}
            },
            {
                "task_key": "edit_status",
                "depends_on": [{"task_key": "Chat_Master_Index"},{"task_key": "Chat_Repo_Index"},{"task_key": "Search_Master_Index"},{"task_key": "Search_Repo_Index"}],
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE"
                },
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "timeout_seconds": 0,
                "email_notifications": {},
                "notification_settings": {
                    "no_alert_for_skipped_runs": False,
                    "no_alert_for_canceled_runs": False,
                    "alert_on_last_attempt": False
                },
                "webhook_notifications": {}
            }
  ],
        "job_clusters": [
            {
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "new_cluster": 
                {
                    "cluster_name": "",
                    "spark_version": "12.2.x-scala2.12",
                    "spark_conf": {
                    "spark.master": "local[*, 4]",
                    "spark.databricks.cluster.profile": "singleNode"
                    },
                    "azure_attributes": {
                        "first_on_demand": 1,
                        "availability": "ON_DEMAND_AZURE",
                        "spot_bid_max_price": -1},
                    "node_type_id": "Standard_D8ads_v5",
                    "driver_node_type_id": "Standard_D8ads_v5",
                    "custom_tags": {
                        "ResourceClass": "SingleNode"},
                    "spark_env_vars": {
                        "PYSPARK_PYTHON": "/databricks/python3/bin/python3"},
                    "enable_elastic_disk": True,
                    "init_scripts": [
                        {
                            "workspace": {
                            "destination": ""
                            }
                        }
                    ],
                    "data_security_mode": "SINGLE_USER",
                    "runtime_engine": "STANDARD",
                    "num_workers": 0
                }
            }
        ],
        "queue": {
            "enabled": True},
        "parameters": [
            {
                "name": "folder_name",
                "default": folder_name
            },
            {
                "name": "job_id",
                "default": "{{job.id}}"
            },
            {
                "name": "site_id",
                "default": site_id
            },
            {
                "name": "workspace",
                "default": workspace
            }
        ],
        "run_as": {
            "user_name": ""
        }
    }


    # API endpoint for creating a job
    create_job_url = f"{workspace_url}/api/2.1/jobs/create"

    # Make the request to create the job
    response = requests.post(
        create_job_url,
        json=job_config,
        headers={"Authorization": f"Bearer {token}"},
    )

    # Check if the job creation was successful
    if response.status_code == 200:
        print("Scheduled job created successfully.")
        job_id = response.json()["job_id"]
        print("Job ID:", job_id)
        return True, job_id
    else:
        print("Failed to create scheduled job.")
        print("Error:", response.text)
        return False, response.text

def update_job(pipeline):
    workspace_url = client.get_secret("DATABRICKS-URL").value
    token = client.get_secret("DATABRICKS-PAT").value
    
    job_id = pipeline['databricksjobid']
    # Get the current job configuration
    get_job_url = f"{workspace_url}/api/2.0/jobs/get?job_id={job_id}"

    # Make the request to get the current job configuration
    response = requests.get(
        get_job_url,
        headers={"Authorization": f"Bearer {token}"},
    )

    # Check if the request to get the job configuration was successful
    if response.status_code == 200:
        current_job_config = response.json()["settings"]
        
        # Update the job configuration as needed
        current_job_config["schedule"]["quartz_cron_expression"] = generate_cron_expression(pipeline['frequency'],pipeline['time'],pipeline['day'])
        
        # API endpoint for updating a job
        update_job_url = f"{workspace_url}/api/2.0/jobs/reset"

        # Make the request to update the job
        update_response = requests.post(
            update_job_url,
            json={"job_id": job_id, "new_settings": current_job_config},
            headers={"Authorization": f"Bearer {token}"},
        )

        # Check if the job update was successful
        if update_response.status_code == 200:
            print("Scheduled job updated successfully.")
        else:
            print("Failed to update scheduled job.")
            print("Error:", update_response.text)
    else:
        print("Failed to get current job configuration.")
        print("Error:", response.text)

def run_job(job_id, notebook_params = False):
    job_run_api = f"{workspace_url}/api/2.1/jobs/run-now"

    response = requests.post(
        job_run_api,
        json={"job_id": job_id} if not notebook_params else {"job_id": job_id,"notebook_params":notebook_params},
        headers={"Authorization": f"Bearer {token}"},
    )

    if response.status_code == 200:
        print("job ran successfully.")
        run_id = response.json()["run_id"]
        print("Run ID:", run_id)
        return True, run_id
    else:
        print("failed to run the job")
        print("Error:", response.text)
        return False, response.text
    
def create_update_job(pipelines,workspace):
    try:
        db = Database('pipelineSettings',workspace)

        for pipeline in pipelines:
            if not pipeline['databricksjobid']:
                status,job_id = create_job(pipeline,workspace)
                status_job_run=run_job(job_id)
                if status:
                    pipeline["databricksjobid"] = job_id
                    pipeline['created_on'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
                    pipeline['modified_on'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
                    pipeline.pop('_id',None)
                    db.insert_single_record(pipeline)
                else:
                    print(f"Failed creation of {pipeline['pipeline_name']}")
            if pipeline['updated']:
                update_job(pipeline)
                id=pipeline['_id']
                pipeline.pop('_id',None)
                pipeline['modified_on']=datetime.now(ist_timezone).strftime('%Y-%m-%d %H:%M:%S')
                db.update_one_record(id,pipeline)
        return True, "Pipeline(s) sucessfully added/updated."
    except Exception as error:
        return False, f"Failed to add/update pipeline(s). Reason: {error}"

def delete_job(id,workspace):
    workspace_url = client.get_secret("DATABRICKS-URL").value
    token = client.get_secret("DATABRICKS-PAT").value
    db = Database('pipelineSettings',workspace)
    status,record = db.fetch_one_record({"_id":id},{'databricksjobid':1,'foldername':1})
    if status:

        background_thread = threading.Thread(target=delete_pipeline_data, args=(record['foldername'],workspace))
        background_thread.start()
        job_id = record['databricksjobid']
        delete_job_url = f"{workspace_url}/api/2.1/jobs/delete"
        response = requests.post(
            delete_job_url,
            json={"job_id":job_id},
            headers={"Authorization": f"Bearer {token}"}
        )
        if response.status_code == 200:
            print("Scheduled job deleted successfully.")
            db.delete_single_record(id)
            return True
        else:
            print("Failed to delete scheduled job.")
            print("Error:", response.text)
            return False

def delete_pipeline_data(folder_name,workspace):
    try:
        db = Database('DocumentCollection',workspace)
        db_central=Database('workspace','central')
        print("Retrieving the record to be deleted.")
        status,results = db.fetch_all_records({"folder_name" : folder_name})
        if status:
            status,records=db_central.fetch_one_record({"workspace_name":workspace},{"_id":0,"upload_job_id":1})
            wokspaceurl=DATABRICKS_URL
            databricks_api_url = f"{wokspaceurl}/api/2.1/jobs/run-now"
            token = PAT_TOKEN
            
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {token}"
            }
            payload = {
                "job_id": records['upload_job_id'],
                "job_parameters": {"type":"pipeline_delete","folder_name":folder_name, "workspace":workspace}
            }

            logging.info("Sending request to Databricks...")
            response = requests.post(databricks_api_url, headers=headers, json=payload) 
    except Exception as error:
        print(f"utility function - delete_data - Error : {error}")
        return False,"Failed to delete file and record."
        
def get_access_token(workspace):
    db_central=Database('workspace','central')
    query=[{'$match': {'workspace_name': workspace,"clientID":{"$not":{"$eq":None}}}},{"$project":{"_id":0,"clientID":1,"tenantID":1,"clientSecret":1}}]
    status_connector,result_connector=db_central.fetch_aggregate(query)
    if status_connector:
        client_id,tenant_id,client_secret = result_connector["clientID"],result_connector["tenantID"],result_connector["clientSecret"]

        url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/token"

        payload = f'grant_type=client_credentials&client_id={client_id}&client_secret={client_secret}&resource=https%3A%2F%2Fgraph.microsoft.com%2F'

        headers = {'Content-Type': 'application/x-www-form-urlencoded'}
        response = requests.request("POST", url, headers=headers, data=payload)

        if response.status_code==200:
            access_token = response.json()['access_token']
            return True,access_token
        
        if response.status_code==401:
            return False,response.json()['error_description']

def get_sharepoint_download_link(access_token,site_id,file_name,workspace):
    url = f"{base_url}/{site_id}/drive/root:{file_name}"
    file_url = ''
    while True:
        payload = ""
        headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {access_token}'
        }
        response = requests.request("GET", url, headers=headers, data=payload)
        response = response.json()
        if '@microsoft.graph.downloadUrl' in response:
            file_url = response['@microsoft.graph.downloadUrl']
            break
        else:
            print("Getting access token.")
            status,access_token = get_access_token(workspace)
            if not status:
                break
    return file_url
    

def generate_link(url,workspace):
    db = Database('DocumentCollection',workspace)
    status,record = db.fetch_one_record({'file_path':url},{'site_id':1,'item_name':1})
    if status:
        status,access_token = get_access_token(workspace)
        if status:
            download_url = get_sharepoint_download_link(access_token,record['site_id'],record['item_name'])
            if download_url:
                return download_url
    return url

def table_extraction(table):
    table_data = []

    for cell in table.cells:
        while len(table_data) <= cell.row_index:
            table_data.append([])

        while len(table_data[cell.row_index]) <= cell.column_index:
            table_data[cell.row_index].append("")

        table_data[cell.row_index][cell.column_index] = cell.content
    table_text = ""
    for row in table_data:
        formatted_row = ','.join(row)
        table_text = table_text+formatted_row+"\n"
    return table_text

def page_to_pdf(doc,pages):
    new_pdf = fitz.open()
    new_pdf.insert_pdf(doc,from_page=pages-1, to_page=pages-1)
    bytes_io = BytesIO()
    new_pdf.save(bytes_io)
    pdf_bytes = bytes_io.getvalue()
    new_pdf.close()
    return pdf_bytes

def doc_intelligence(file):
    #print("file ",file)
    poller = document_analysis_client.begin_analyze_document('prebuilt-read',document= file)
    result = poller.result()
    doc_pages = result.pages
    tables = result.tables
    #print("tables ",tables)
    for page in doc_pages:
        lines_content = []
        for line in page.lines:
            lines_content.append(line.content)
        #print("line ",line)
        page_text = (" ".join(lines_content)) 
        for table in tables:
            if table.bounding_regions and table.bounding_regions[0].page_number==page.page_number:
                page_text += table_extraction(table)
        #print(" page_text",page_text)
    return page_text

def get_pdf_text(doc,page,page_number):

    page_text = ""
    table = page.find_tables()
    image = page.get_images()
    page_text = page.get_text()
    if len(table.tables):
        #print("tables",page_number)
        pdf_bytes = page_to_pdf(doc, page_number)
        print(type(pdf_bytes))
        print("getting table result")
        try:
            result = doc_intelligence(pdf_bytes)
            print(type(result))
            return result
        except Exception as e:
            print("Exception in table",e)
    elif len(image) and len(page_text) == 0:
        #print("imaged and text",page_number)
        pdf_bytes = page_to_pdf(doc, page_number)
        #print(type(pdf_bytes))
        try:
            result = doc_intelligence(pdf_bytes)
            #print(type(result))
            return result
        except Exception as e:
            print("Exception in image,e")
    else:
        print("else",page_number)
        for image_index, img in enumerate(image, start=1):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                img = Image(src=BytesIO(image_bytes))
                img_tables = img.extract_tables()
                if len(img_tables):
                    pdf_bytes = page_to_pdf(doc, page_number)
                    print(type(pdf_bytes))
                    try:
                        result = doc_intelligence(pdf_bytes)
                        print(type(result))
                        return result
                    except Exception as e:
                        print("Exception is ",e)
                else:
                    print("test")
                    try:
                        status, text = ocr_img2text(image_bytes)
                        if status:
                            text = re.sub(' +', ' ', text)
                            page_text += text
                    except Exception as e:
                        print("Exception",e)
                    
            except Exception as e:
                pass
    #await get_page_data(page_number,)
    return page_text

def process_all_pages(doc):
    arr_page = []
    pages = [page for page in doc]
    page_number=0
    for page in pages:
        future_result = get_pdf_text(doc,page, page_number + 1)
        result = future_result
        page_number+=1
        #print(result)
        arr_page.append(result)

    return arr_page

def get_pdf_fulltext(file):
    page_text = ""
    doc = fitz.open(stream=file)
    page_meta = process_all_pages(doc)
    for page in page_meta:
        page_text+=page
    return True,page_text


import pathlib
def continuous_recognize_text_from_audio(file):
    
    # This example requires environment variables named "SPEECH_KEY" and "SPEECH_REGION"
    speech_config = speechsdk.SpeechConfig(subscription=client.get_secret("AZURE-SPEECH-KEY").value, region="centralindia")
    # speech_config.speech_recognition_language = "hi-IN"

    # Use the local file as the audio input
    audio_config = speechsdk.audio.AudioConfig(filename=file)

    from_languages = ["en-IN","mr-IN","hi-IN"]

    auto_detect_config = AutoDetectSourceLanguageConfig(from_languages)
 
    # recognizer = speechsdk.SpeechRecognizer(speech_config=speech_config, audio_config=audio_config)
    recognizer = SpeechRecognizer(speech_config=speech_config, audio_config=audio_config, auto_detect_source_language_config=auto_detect_config)
    done = False
    recognized_text = ""

    def stop_cb(evt):
        nonlocal done
        done = True
 
    def on_continuous_recognize(evt):
        nonlocal recognized_text
        if evt.result.reason == speechsdk.ResultReason.RecognizedSpeech:
            recognized_text += evt.result.text
        elif evt.result.reason == speechsdk.ResultReason.NoMatch:
            print("No speech could be recognized: {}".format(evt.result.no_match_details))
              # Handle NoMatch as needed
        elif evt.result.reason == speechsdk.ResultReason.Canceled:
            pass  # Handle Canceled as needed
 
    # recognizer.recognizing.connect(on_continuous_recognize)
    recognizer.recognized.connect(on_continuous_recognize)
    recognizer.session_started.connect(lambda evt: print('SESSION STARTED: {}'.format(evt)))
    recognizer.session_stopped.connect(lambda evt: print('SESSION STOPPED {}'.format(evt)))
    recognizer.canceled.connect(lambda evt: print('CANCELED {}'.format(evt)))
    recognizer.session_stopped.connect(stop_cb)
    recognizer.canceled.connect(stop_cb)
 
    recognizer.start_continuous_recognition()
    time.sleep(20)
    while not done:
        # time.sleep(10)
        pass
 
    recognizer.stop_continuous_recognition()
    return recognized_text

def translate_english_to_hindi(text, lang):
    ch_text = [text]
    subscription_key = client.get_secret("AZURE-TRANSLATION-KEY").value
    endpoint = ""
    location = "global"
    path = '/translate'
    constructed_url = endpoint + path
    params = {
        'api-version': '3.0',
        'to': lang
    }
    constructed_url = endpoint + path
    headers = {
        'Ocp-Apim-Subscription-Key': subscription_key,
        'Ocp-Apim-Subscription-Region': location,
        'Content-type': 'application/json',
        'X-ClientTraceId': str(uuid.uuid4())
    }
    cnv = []
    for i in range(0, len(ch_text)):
        # You can pass more than one object in body.
        body = [{
            'text': str(ch_text[i])
        }]
        request = requests.post(constructed_url, params=params, headers=headers, json=body)
        response = request.json()
        cnv.append(response[0]['translations'][0]['text'])

    # Concatenate the translations into a single string
    result_string = ' '.join(cnv)
    return result_string

def audio_to_text(file):
    # audio_data = {}
    parent = str(pathlib.Path().absolute())
    local_file_path = parent+f'/test_{datetime.now().strftime("%Y-%m-%d %H:%M:%S")}.wav'
    with open(local_file_path, "wb") as audio_file:
        audio_file.write(file)
    full_text = continuous_recognize_text_from_audio(local_file_path)
    # print(full_text)
    print("File transcribed successfully.")
    translated_text = translate_english_to_hindi(full_text,'en-IN')
    print("File translated successfully.")
    clean_text_result = clean_text_audio(translated_text)
    print("Text cleaned sucessfully.")
    os.remove(local_file_path)
    # os.unlink(local_file_path)
    print("File deleted sucessfully.")
    return True,clean_text_result

def clean_text_audio(input_text):
    try:
        if not input_text.strip():
            return ""
        # value = 0
        AZURE_OPENAI_CHATGPT_DEPLOYMENT = "gpt-4-turbo"

        openai_client = openai.AzureOpenAI(
            azure_endpoint = client.get_secret("AZURE-OPENAI-ENDPOINT-IN").value,
            api_version="2023-09-01-preview",
            api_key = client.get_secret('OPENAI-API-KEY-IN').value

        )

        messages = [
            {"role":"system","content":"""Revise the text to enhance readability by eliminating redundancies and correcting spelling errors. Return an empty string if Audio trancription contain no text. Make sure all the context that can be extracted from the input is covered:"""},
            {"role": "user", "content": f"Text:\n\n{input_text}\n\nCleaned Text:\n-"}]
        response = openai_client.chat.completions.create(
            model=AZURE_OPENAI_CHATGPT_DEPLOYMENT,
            messages=messages,
            temperature=0,
            max_tokens=2000,
            timeout=50
        )
        text = response.choices[0].message.content
        print("Cleaned_text: ",text)
        if text:
            return text.strip()
        return ""  
    except Exception as error:
        print(f"Failed->summary_generation. Error {error}")
        return ""

def get_file_text(file,ext):
    full_text = ""
    status = False
    if ext in ['pptx','ppt']:
        presentation = slides.Presentation(file)
        pdf_file = BytesIO()
        presentation.save(pdf_file, slides.export.SaveFormat.PDF)
        status,full_text = get_pdf_fulltext(pdf_file)
        pdf_file.flush()
        pdf_file.seek(0)
    elif ext == 'pdf':
        status,full_text = get_pdf_fulltext(file)
    elif ext in ['doc','docx']:
        doc = aw.Document(file)
        pdf_file = BytesIO()
        doc.save(pdf_file,aw.SaveFormat.PDF)
        status,full_text = get_pdf_fulltext(pdf_file)
        pdf_file.flush()
        pdf_file.seek(0)
    elif ext in ['wav']:
        status,full_text = audio_to_text(file.getbuffer())
        print(full_text)
    return status,full_text

def get_llm_response(instructions,data):
    AZURE_OPENAI_CHATGPT_DEPLOYMENT = "gpt-4-turbo"
    openai.azure_endpoint = client.get_secret('AZURE-OPENAI-ENDPOINT-IN').value
    openai.api_key = client.get_secret('OPENAI-API-KEY-IN').value
    
    functions= [  
    {
        "name": "Application",
        "description": "Perform the following task on the given Data",
        "parameters": {
            "type": "object",
            "properties": {
                f"{instruction['title']}": {
                    "type": "string",
                    "description": f"{instruction['description']}"
                }
            for instruction in instructions
            },
            "required": [f"{instruction['title']}" for instruction in instructions]
        }
    }
    ]
    msg = [
        {"role": "user", 
         "content": data}]
    completion = openai.chat.completions.create(
        model=AZURE_OPENAI_CHATGPT_DEPLOYMENT,
        messages=msg,
        functions=functions,
        function_call='auto',
        temperature=0,
        max_tokens=512,
        timeout=100)
    result = completion.choices[0].message
    result = eval(result.function_call.arguments)
    result = [value for key,value in result.items()]
    results = [{'title':instruction['title'],'description':instruction['description'],'output':output} for instruction,output in zip(instructions,result)]
    return results


def get_output(file,instructions):
    ext = file["file_name"].split(".")[-1]
 
    status,file_text = get_file_text(file['data'],ext)
    if status:
        print(file_text)
        result = get_llm_response(instructions,file_text)
            
    return result


def test_user_application(test_data,files,workspace):
    try:
        db = Database('TestApplication',workspace)
        print(test_data)
        instructions = test_data['instructions']
        data = []
        for file in files:
            print(file['file_name'])
            result = {
                    "file_name":file['file_name']   
                }
            output = get_output(file,instructions)
            result.update({"results": output})
            data.append(result)
        test_data.update({"files":[file['file_name'] for file in files]})
        test_data.update({"result":data})
        if not test_data['test_id']:
            test_data.pop("test_id")
            status,id = db.insert_single_record(test_data)
            test_data['test_id'] = id
            test_data.pop("_id")
        else:
            test_record = test_data
            id = test_record.pop("test_id")
            test_record['_id'] = id
            status,id = db.update_one_record(id,test_record)
        if 'result' in test_data:
            response = {
                "message":"Files testing successfull",
                "data": test_data['result']
            }
            return True,response
    except Exception as error:
        print(error)
        return False,{}

def create_application_pipeline(application_id,pipeline,workspace):
    pipeline_name = pipeline['pipeline_name']
    pipeline_id = str(pipeline['pipeline_id'])
    application_id = str(application_id)
    job_config = {
        "name": f"{workspace}_Application_{pipeline_name}",
        "job_clusters": [
            {
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}",
                "new_cluster": {
                    "cluster_name": "",
                    "spark_version": "12.2.x-scala2.12",
                    "spark_conf": {
                        "spark.master": "local[*, 4]",
                        "spark.databricks.cluster.profile": "singleNode"
                    },
                    "azure_attributes": {
                        "first_on_demand": 1,
                        "availability": "ON_DEMAND_AZURE",
                        "spot_bid_max_price": -1.0
                    },
                    "node_type_id": "Standard_D4s_v3",
                    "driver_node_type_id": "Standard_D4s_v3",
                    "custom_tags": {
                        "ResourceClass": "SingleNode"
                    },
                    "enable_elastic_disk": True,
                    "init_scripts": [
                        {
                            "workspace": {
                                "destination": ""
                            }
                        }
                    ],
                    "data_security_mode": "SINGLE_USER",
                    "runtime_engine": "STANDARD",
                    "num_workers": 0
                }
            }
        ],
        "email_notifications": {},
        "tasks":[
            {
                "task_key": f"AIXponentApplication_{workspace}",
                "run_if": "ALL_SUCCESS",
                "notebook_task": {
                    "notebook_path": "",
                    "source": "WORKSPACE",
                    "base_parameters": {
                        "workspace":workspace,
                        "pipeline_id": pipeline_id,
                        "application_id": application_id,
                        "job_id":"{{job_id}}"}},
                "job_cluster_key": f"GPTXponentJobCluster_{workspace}"
        }],
        "schedule": {
            "quartz_cron_expression": generate_cron_expression(pipeline['frequency'].title(),pipeline['time'],pipeline['day']),  # Every day at midnight
            "timezone_id": "Asia/Kolkata",  # Set your desired timezone
        },
        "max_concurrent_runs": 5,
        "format": "MULTI_TASK",
        "queue": {
            "enabled": True
        },
        "run_as": {
            "user_name": ""
        }
    }


    # API endpoint for creating a job
    create_job_url = f"{DATABRICKS_URL}/api/2.0/jobs/create"

    # Make the request to create the job
    response = requests.post(
        create_job_url,
        json=job_config,
        headers={"Authorization": f"Bearer {PAT_TOKEN}"},
    )

    # Check if the job creation was successful
    if response.status_code == 200:
        print("Scheduled job created successfully.")
        job_id = response.json()["job_id"]
        print("Job ID:", job_id)
        return True, job_id
    else:
        print("Failed to create scheduled job.")
        print("Error:", response.text)
        return False, response.text


def create_update_application_pipelines(application_id,pipelines,workspace,username):
    try:
        db = Database('ApplicationCollection',workspace)
        status,record = db.fetch_one_record({"_id":application_id},{'pipelines':1})
        for pipeline,old_pipeline in zip(pipelines,record['pipelines']):
            pipeline['pipeline_id']= ObjectId(pipeline['pipeline_id'])
            if not pipeline['databricksjobid']:
                print('Started pipeline creation.')
                pipeline['pipeline_id']= ObjectId()
                status,job_id = create_application_pipeline(application_id,pipeline,workspace)
                print('Pipeline creation completed.')
                status_job_run=run_job(job_id)
                print('Pipeline job running.')
                if status:
                    pipeline["databricksjobid"] = job_id
                    pipeline['instruction_updated'] = False
                    pipeline['created_by'] = username
                    pipeline['created_at'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
                    pipeline['modified_by'] = username
                    pipeline['modified_at'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
                else:
                    print(f"Failed creation of {pipeline['pipeline_name']}")
            if pipeline['updated']:
                print('Pipeline job being updated.')
                update_job(pipeline)
                print('Pipeline job updated.')
                pipeline['update'] = False
                pipeline['instruction_updated'] = True if pipeline['instructions'] != old_pipeline['instructions'] else False
                print('Pipeline creation completed.')
                status_job_run=run_job(pipeline['databricksjobid'])
                pipeline['modified_by'] = username
                pipeline['modified_at'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
        return True, pipelines
    except Exception as error:
        print(f"Failed to add/update pipeline(s). Reason: {error}")
        return False, f"Failed to add/update pipeline(s). Reason: {error}"

def data_to_cosmos(db_type, data, workspace, username):
    try:
        db = Database("ApplicationCollection", workspace)
        
        if db_type == 'insert':
            data.pop('application_id')
            status, id = db.insert_single_record(data)
            print("inserted id:",id)
            status,pipelines = create_update_application_pipelines(id,data['pipelines'],workspace,username)
            data['pipelines'] = pipelines if status else data['pipelines']
            data['created_by'] = username
            data['created_at'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
            data['modified_by'] = username
            data['modified_at'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
                
        elif db_type == 'update':
            print('Started application update.')
            # Assuming '_id' field is a string representation of ObjectId
            id = data.pop('application_id')
            print('Acquired application_id.')
            status,pipelines = create_update_application_pipelines(id,data['pipelines'],workspace,username)
            print('Application updation done.')
            data['_id'] = id
            data['pipelines'] = pipelines if status else data['pipelines']
            data['modified_by'] = username
            data['modified_at'] = datetime.now(ist_timezone).strftime("%Y-%m-%d %H:%M:%S")
            print("Application update done")

        status, msg = db.update_one_record(id,data)
        if status:
            print("Application creation/update done")
            return status, id
        else:
            return status, msg
    except Exception as e:
        print(f"Failed to insert record: {str(e)}")
        return False, f"Failed to insert record: {str(e)}"

def delete_application_data(job_id, workspace):
    try:
        db = Database('ApplicationDocuments',workspace)
        status,message = db.delete_all_record({'job_id':str(job_id)})
    except Exception as error:
        print("Failed to delete application documents.")

def delete_application_job(job_id,workspace):
    workspace_url = client.get_secret("DATABRICKS-URL1").value
    token = client.get_secret("DATABRICKS-PAT1").value
    status=True
    if status:

        background_thread = threading.Thread(target=delete_application_data, args=(job_id,workspace))
        background_thread.start()
        delete_job_url = f"{workspace_url}/api/2.1/jobs/delete"
        response = requests.post(
            delete_job_url,
            json={"job_id":str(job_id)},
            headers={"Authorization": f"Bearer {token}"}
        )
        if response.status_code == 200:
            print("Scheduled job deleted successfully.")
            return True
        else:
            print("Failed to delete scheduled job.")
            print("Error:", response.text)
            return False

def delete_app_pipeline(application_id,pipeline_id = None,workspace = None):
    try:
        db = Database('ApplicationCollection',workspace)
        if pipeline_id:
            status,records = db.fetch_aggregate([
                {
                    '$match': {
                        '_id': ObjectId(application_id)
                    }
                }, {
                    '$unwind': {
                        'path': '$pipelines'
                    }
                }, {
                    '$match': {
                        'pipelines.pipeline_id': ObjectId(pipeline_id)
                    }
                }, {
                    '$replaceRoot': {
                        'newRoot': '$pipelines'
                    }
                }
            ])
        else:
            status,record = db.fetch_one_record({'_id':application_id},{'_id':0,'pipelines':1})
            if status:
                records = record['pipelines']
        if status:
            for record in records:
                job_id = record['databricksjobid']
                status = delete_application_job(job_id,workspace)
                if status:
                    status,message = db.append_data({'_id':application_id},{'$pull':{'pipelines':{'pipeline_id':record['pipeline_id']}}})
                    if status:
                        return True, "Sucessfully deleted pipeline"
                    else:
                        return False, message
        else:
            return False, "Failed to fetch pipeline details"    
    except Exception as error:
        print(f"Failed to delete pipeline. Error: {error}")
        return False, "Failed to delete pipeline."

def delete_application(application_id, workspace):
    try:
        status,message = delete_app_pipeline(application_id=application_id,workspace=workspace)
        if status:
            db = Database('ApplicationCollection',workspace)
            status,message = db.delete_all_record({'_id':application_id})
            if status:
                return True, "Sucessfully deleted application"
            else:
                return False, message
        else:
            return False, "Failed to delete pipelines."
    except Exception as error:
        print(f"Failed to delete application. Error: {error}")
        return False, "Failed to delete application."


def application_output(application_id, workspace, app_base_url):
    try:
        final_result = {}
        db = Database('ApplicationCollection',workspace)
        status,application = db.fetch_one_record({'_id':application_id})
        if not status:
            return False, "Application not present."
        final_result = {
            "workspace_id":workspace,
            "application_id": str(application['_id']),
            "application_name":application['application_name'],
            "created_at":application['created_at'],
            "created_by" : application['created_by'],
            "modified_at":application['modified_at'],
            "modified_by" : application['modified_by'],
            "output_api" : f"{app_base_url}tenant/{workspace}/applications/{str(application['_id'])}/output-schema"
        }
        print("pipelines")
        pipelines= []
        for pnum,pipeline in enumerate(application['pipelines']):
            pipeline_data = {}
            pipeline_data['key'] = f"0-{pnum}"
            pipeline_data['count'] = f"{pnum+1}"
            pipeline_data['label'] = pipeline['pipeline_name']
            pipeline_data['type'] = "pipeline"
            pipeline_data['children'] = [ 
              {
                "key": f"0-{pnum}-0",
                "label": "Created At",
                "value": f"{pipeline['created_at']}",
                "type": "timestamp"
              },
              {
                "key": f"0-{pnum}-1",
                "label": "Created By",
                "value": f"{pipeline['created_by']}",
                "type": "user"
              },
              {
                "key": f"0-{pnum}-2",
                "label": "Modified At",
                "value": f"{pipeline['modified_at']}",
                "type": "timestamp"
              },
              {
                "key": f"0-{pnum}-3",
                "label": "Modified By",
                "value": f"{pipeline['modified_by']}",
                "type": "user"
              },
              {
                "key": f"0-{pnum}-4",
                "label": "Last Run At",
                "value": f"{pipeline['last_run_at']}",
                "type": "timestamp"
              },
              {
                "key": f"0-{pnum}-5",
                "label": "Last Run Completed At",
                "value": f"{pipeline['last_run_completed_at']}",
                "type": "timestamp"
              },
              {
                "key": f"0-{pnum}-6",
                "label": "Source Folders"
              }]
            folders = []
            print("folder")
            for fnum,folder in enumerate(pipeline['source_details']):
                folder_data ={
                    "key" : f"0-{pnum}-6-{fnum}",
                    "label" : folder['folder_name'].split("/")[-1]
                }
                db = Database('ApplicationDocuments',workspace)
                status, files = db.fetch_all_records({'job_id':str(pipeline['databricksjobid']),'folder_name':str(folder['folder_name']),"edit_flag":1},{'file_name':1,'file_path':1,'results':1})
                files_list = []
                print("files")
                if status:
                    for file_num,file in enumerate(files):
                        files_dict = {
                            "key": f"0-{pnum}-6-{fnum}-{file_num}",
                            "label":file['file_name'],
                            "children": [
                            {
                                "key": f"0-{pnum}-6-{fnum}-{file_num}-0",
                                "label": "Instructions",
                                "children": []
                                
                                }
                                ]
                        }
                        print("results")
                        results = []
                        
                        for result_num,result in enumerate(file['results']):
                            result_dict = {
                                "key": f"0-{pnum}-6-{fnum}-{file_num}-0-{result_num}",
                                "label":f"Instruction {result['instruction_number']}"
                            }
                            # result_items = []
                            result.pop('instruction_number',None)
                            result_items = {
                                "key":f"0-{pnum}-6-{fnum}-{file_num}-0-{result_num}-0",
                                **result,
                                "type":"instruction"
                            }

                            result_dict['children'] = [result_items]
                            results.append(result_dict)
                        files_dict['children'][0]['children'] = results
                        files_list.append(files_dict)
                folder_data['children'] = files_list
                folders.append(folder_data)
            pipeline_data['children'][-1]['children'] = folders
            pipelines.append(pipeline_data)
        final_result['pipelines'] = [{
            "key": "0",
            "label": "Pipelines",
            "children": pipelines
            }]
        return True, final_result
            
    except Exception as Error:
        print(f"Failed to retrieve application output. Error: {Error}")
        return False, "Failed to retrieve application output."

def application_output_external(application_id, workspace, app_base_url):
    try:
        final_result = {}
        db = Database('ApplicationCollection',workspace)
        status,application = db.fetch_one_record({'_id':application_id})
        if not status:
            return False, "Application not present."
        final_result = {
            "workspace_id":workspace,
            "application_id": str(application['_id']),
            "application_name":application['application_name'],
            "created_at":application['created_at'],
            "created_by" : application['created_by'],
            "modified_at":application['modified_at'],
            "modified_by" : application['modified_by'],
            "output_api" : f"{app_base_url}applications/{str(application['_id'])}/output"
        }
        print("pipelines")
        pipelines= []
        for pipeline in application['pipelines']:
            pipeline_data = {}
            pipeline_data['pipeline_name'] = pipeline['pipeline_name']
            pipeline_data["created_at"] = pipeline['created_at']
            pipeline_data["created_by"] = pipeline['created_by']
            pipeline_data["modified_at"] = pipeline['modified_at']
            pipeline_data["modified_by"] = pipeline['modified_by']
            pipeline_data["last_run_at"] = pipeline['last_run_at']
            pipeline_data["last_run_completed_at"] = pipeline['last_run_completed_at']
            folders = []
            print("folder")
            for fnum,folder in enumerate(pipeline['source_details']):
                folder_data ={
                    "folder_name" : folder['folder_name']
                }
                db = Database('ApplicationDocuments',workspace)
                status, files = db.fetch_all_records({'job_id':str(pipeline['databricksjobid']),'folder_name':str(folder['folder_name']),"edit_flag":1},{'file_name':1,'file_path':1,'results':1})
                files_list = []
                print("files")
                if status:
                    for file_num,file in enumerate(files):
                        files_dict = {
                            "file_name":file['file_name'],
                            "file_path":file['file_path'],
                            "results": file['results']
                        }
                        files_list.append(files_dict)
                folder_data['files'] = files_list
                folders.append(folder_data)
            pipeline_data['folders'] = folders
            pipelines.append(pipeline_data)
        final_result['pipelines'] = pipelines
        return True, final_result
            
    except Exception as Error:
        print(f"Failed to retrieve application output. Error: {Error}")
        return False, "Failed to retrieve application output."


def file_process1(id,workspace):
    db=Database("workspace","central")
    db_database=Database("DocumentCollection",workspace)
    status,records=db.fetch_one_record({"workspace_name":workspace},{"_id":0,"upload_job_id":1})
    # print(records)
    if status and "upload_job_id" in records:
        wokspaceurl=DATABRICKS_URL
        databricks_api_url = f"{wokspaceurl}/api/2.1/jobs/run-now"
        token = PAT_TOKEN
        
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {token}"
        }
        payload = {
            "job_id": records['upload_job_id'],
            "job_parameters": {"type":"insert","doc_id":id, "workspace":workspace}
        }

        logging.info("Sending request to Databricks...")
        response = requests.post(databricks_api_url, headers=headers, json=payload)
        if response.status_code == 200:
            try:
                notebook_response = response.json()
                run_id = notebook_response.get("run_id")
                status, message = db_database.update_one_record(id, {"run_id":run_id})
                url = f"{wokspaceurl}/api/2.1/jobs/runs/get?run_id={run_id}"
                payload = ""
                headers = {"Authorization": f"Bearer {token}"}
        
                response = requests.request("GET", url, headers=headers, data=payload)
        
                if response.status_code==200:
                    notebook_response = response.json()
                    try:
                        task_id= list(filter(lambda x: x["task_key"] == "File_Process1", notebook_response["tasks"]))[0]['run_id']
                    except Exception as e:
                        print(e)
                    url = f"{wokspaceurl}/api/2.1/jobs/runs/get-output?run_id={task_id}"
                    status = notebook_response["state"]["life_cycle_state"]
        
                    while status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                        response = requests.request("GET", url, headers=headers, data=payload)
                        if response.status_code==200:
                            notebook_response = response.json()
                            status = notebook_response['metadata']["state"]["life_cycle_state"]

                        else:
                            print(json.dumps({"notebook_response": "Error in fetching id"}))
                        if status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                            time.sleep(5)
                    data=notebook_response
                    # print(data)
                    job_status=notebook_response['metadata']["state"]["result_state"]
                    if job_status=='SUCCESS':
                        return True, "Successfully processed file"
                    else:
                        return False, "Error In File Process"
                return False, "Unable to process file"
            except Exception as e:
                return False, f"Error: full_process->{e}"
    else:
        return False, "Job Not Available"

def file_process2(id,workspace):
    db=Database("DocumentCollection",workspace)
    status,records=db.fetch_one_record({"_id":id},{"_id":0,"run_id":1})
    print(records)
    if status and "run_id" in records:
        wokspaceurl=DATABRICKS_URL
        
        run_id = records.get("run_id")
        
        url = f"{wokspaceurl}/api/2.1/jobs/runs/get?run_id={run_id}"
        payload = ""
        headers = {"Authorization": f"Bearer {token}"}

        response = requests.request("GET", url, headers=headers, data=payload)

        if response.status_code==200:
            notebook_response = response.json()
            task_id= list(filter(lambda x: x["task_key"] == "file_process2", notebook_response["tasks"]))[0]['run_id']
            url = f"{wokspaceurl}/api/2.1/jobs/runs/get-output?run_id={task_id}"
            status = notebook_response["state"]["life_cycle_state"]

            while status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                response = requests.request("GET", url, headers=headers, data=payload)
                if response.status_code==200:
                    notebook_response = response.json()
                    status = notebook_response['metadata']["state"]["life_cycle_state"]

                else:
                    print(json.dumps({"notebook_response": "Error in fetching id"}))
                if status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                    time.sleep(5)
            data=notebook_response
            # print(data)
            job_status=notebook_response['metadata']["state"]["result_state"]
            if job_status=='SUCCESS':
                return True, "Successfully processed file"
            else:
                return False, "Error In File Process2"
        return False, "Unable to process file"
    else:
        return False, "Process1 Has Not Started"
    

def file_preview(id,workspace):
    while True:
        db=Database("DocumentCollection",workspace)
        status,records=db.fetch_one_record({"_id":id},{"_id":0,"run_id":1})
        if status and "run_id" in records:
            wokspaceurl=DATABRICKS_URL
            
            run_id = records.get("run_id")
            
            url = f"{wokspaceurl}/api/2.1/jobs/runs/get?run_id={run_id}"
            payload = ""
            headers = {"Authorization": f"Bearer {token}"}

            response = requests.request("GET", url, headers=headers, data=payload)

            if response.status_code==200:
                notebook_response = response.json()
                try:
                    task_id= list(filter(lambda x: x["task_key"] == "file_preview", notebook_response["tasks"]))[0]['run_id']
                except Exception as e:
                    print(e)
                url = f"{wokspaceurl}/api/2.1/jobs/runs/get-output?run_id={task_id}"
                status = notebook_response["state"]["life_cycle_state"]

                while status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                    response = requests.request("GET", url, headers=headers, data=payload)
                    if response.status_code==200:
                        notebook_response = response.json()
                        status = notebook_response['metadata']["state"]["life_cycle_state"]

                    else:
                        print(json.dumps({"notebook_response": "Error in fetching id"}))
                    if status=="RUNNING" or status=="PENDING" or status=="BLOCKED":
                        time.sleep(5)
                data=notebook_response
                # print(data)
                job_status=notebook_response['metadata']["state"]["result_state"]
                if job_status=='SUCCESS':
                    return True, "Successfully processed file"
                else:
                    return False, "Error In File preview"
            return False, "Unable to process file"
        else:
            time.sleep(2)
            continue

def clean_query_cache():
    # db = Database('ChatCollection',workspace)
    db = Database('ChatCollection')
    today = datetime.today().date()
    previous_date = today - timedelta(days=1)
    previous_date_formatted = previous_date.strftime('%d-%m-%Y')
    status,records = db.fetch_all_records({"created_on": {"$regex": f"^{previous_date_formatted}"}},{"_id":1})
    ids = []
    if status:
        ids = [str(record['_id']) for record in records]
    params ={
        "ids":str(ids),
        "workspace":"pidilite-pd"
    }
    run_job(0123456789,params)